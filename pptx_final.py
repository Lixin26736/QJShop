from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs=Presentation()
prs.slide_width=Inches(13.333)
prs.slide_height=Inches(7.5)

P=RGBColor(0x25,0x63,0xEB);D=RGBColor(0x1A,0x36,0x5D);A=RGBColor(0xF5,0x9E,0x0B)
W=RGBColor(0xFF,0xFF,0xFF);B=RGBColor(0x00,0x00,0x00);G=RGBColor(0x64,0x74,0x8B)
L=RGBColor(0xF1,0xF5,0xF9);R=RGBColor(0xEF,0x44,0x44);E=RGBColor(0x10,0xB9,0x81)
BL=RGBColor(0xDB,0xEA,0xFE);O=RGBColor(0xF9,0x73,0x16);CB=RGBColor(0x1E,0x29,0x3B)
CT=RGBColor(0xE2,0xE8,0xF0);W2=RGBColor(0xBF,0xDB,0xFE)

def bl():return prs.slides.add_slide(prs.slide_layouts[6])
def T(s,l,t,w,h):return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def S(tf,x,fn='Microsoft YaHei',fs=18,bo=False,c=B,al=PP_ALIGN.LEFT):
    tf.word_wrap=1;p=tf.paragraphs[0];p.text=x;p.font.name=fn;p.font.size=Pt(fs)
    p.font.bold=bo;p.font.color.rgb=c;p.alignment=al;return p
def P2(tf,x,fn='Microsoft YaHei',fs=14,bo=False,c=G,al=PP_ALIGN.LEFT,sb=4):
    p=tf.add_paragraph();p.text=x;p.font.name=fn;p.font.size=Pt(fs);p.font.bold=bo
    p.font.color.rgb=c;p.alignment=al;p.space_before=Pt(sb);return p
def X(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background();return sh
def H(s,x,u=None):
    X(s,0,0,13.333,0.06,A);X(s,0,0,13.333,1.15,P)
    tt=T(s,0.7,0.2,11.5,0.7);S(tt,x,fs=30,bo=1,c=W)
    if u:P2(tt,u,fs=14,c=W2,sb=4)
def C(s,lx,ly,w,h,x,y,tc=P):
    X(s,lx,ly,w,h,L);X(s,lx,ly,w,0.42,tc)
    tt=T(s,lx+0.15,ly+0.04,w-0.3,0.35);S(tt,x,fs=14,bo=1,c=W)
    t2=T(s,lx+0.15,ly+0.55,w-0.3,h-0.7);lns=y.split('\n')
    S(t2,lns[0],fs=12,c=D)
    for ln in lns[1:]:P2(t2,ln,fs=11,c=G,sb=3)
def SB(s,lx,ly,w,h,n,lb,co=P):
    X(s,lx,ly,w,h,co);tt=T(s,lx+0.1,ly+h*0.1,w-0.2,h*0.5)
    S(tt,str(n),fs=24,bo=1,c=W,al=PP_ALIGN.CENTER)
    P2(tt,lb,fs=10,c=W2,al=PP_ALIGN.CENTER,sb=4)
def K(s,lx,ly,w,h,x,ls,tc=E):
    X(s,lx,ly,w,h,CB);tt=T(s,lx+0.15,ly+0.08,w-0.3,h-0.2)
    S(tt,x,fs=11,bo=1,c=tc)
    for ln in ls:P2(tt,ln,fs=9,c=CT,sb=3)

# ===== SLIDE 1: COVER =====
s=bl();X(s,0,0,13.333,7.5,P);X(s,0,0,13.333,0.08,A);X(s,0,7.42,13.333,0.08,A)
tt=T(s,1.5,1.2,10.3,1.8);S(tt,'"柒玖商店"',fs=50,bo=1,c=W,al=PP_ALIGN.CENTER)
P2(tt,'基于SpringBoot多前端多平台的商城系统',fs=26,c=W2,al=PP_ALIGN.CENTER,sb=12)
X(s,3.5,3.3,6.3,0.01,W)
tt=T(s,4,3.8,5.5,3);S(tt,'毕业设计答辩',fs=22,bo=1,c=W,al=PP_ALIGN.CENTER)
for ln,sz in [('姓   名：李  鑫',18),('学   号：202330302152',18),('专   业：软件技术  |  班 级：2307',18),('指导教师：彭德宇',18),('2026年5月',16)]:
    P2(tt,ln,fs=sz,c=W if sz>16 else W2,al=PP_ALIGN.CENTER,sb=14)

# ===== SLIDE 2: TOC =====
s=bl();X(s,0,0,0.08,7.5,P);X(s,0,0,13.333,0.06,A)
tt=T(s,0.8,0.4,11,0.8);S(tt,'答辩大纲',fs=36,bo=1,c=D)
toc=[('01','项目背景与选题意义','电子商务发展趋势、市场规模与课题价值'),
    ('02','技术架构概述','前后端技术选型、系统架构与部署方案'),
    ('03','需求分析与功能模块','用户端10模块 + 管理端10模块功能矩阵'),
    ('04','数据库设计','12张核心数据表结构、ER关系与设计特点'),
    ('05','功能演示—用户端(上)','首页/搜索/商品详情(双端)/购物车/结算'),
    ('06','功能演示—用户端(下)','订单管理/收货地址/个人中心/AI客服'),
    ('07','功能演示—管理端(上)','仪表盘/用户管理/分类管理/商品管理'),
    ('08','功能演示—管理端(下)','订单管理/Banner管理/评价管理/系统设置'),
    ('09','AI智能客服(核心亮点)','DeepSeek集成、RAG简化实现、关联词推荐'),
    ('10','作品测试与项目总结','37测试用例、8个技术问题、成果与展望')]
y=1.6
for n,x,d in toc:
    X(s,0.5,y+0.05,0.4,0.4,P)
    tt=T(s,0.53,y+0.07,0.35,0.35);S(tt,n,fs=12,bo=1,c=W,al=PP_ALIGN.CENTER)
    tt=T(s,1.1,y,9,0.3);S(tt,x,fs=16,bo=1,c=D)
    P2(tt,d,fs=11,c=G,sb=2)
    y+=0.57

# ===== SLIDE 3: Background 1 =====
s=bl();H(s,'一、项目背景 — 电商发展现状')
SB(s,0.3,1.6,3.8,1.2,'15万亿+','2025年中国网络零售市场规模',P)
SB(s,4.3,1.6,3.8,1.2,'30%+','占社会消费品零售总额比重',D)
SB(s,8.3,1.6,3.8,1.2,'7×24h','电商突破时空限制全天候服务',E)
tt=T(s,0.5,3.2,6,3.5);S(tt,'行业趋势',fs=20,bo=1,c=D)
for ln in ['移动支付普及率超90%，线上消费习惯已全面养成','AI大模型(DeepSeek等)为电商智能化带来新可能','鸿蒙生态快速扩展，多终端适配成为刚需','中小商家数字化转型需求旺盛但技术门槛高','轻量化、可定制的电商解决方案市场空间巨大']:
    P2(tt,f'• {ln}',fs=14,c=D,sb=12)
K(s,7,3.2,5.5,3.5,'中国电商市场规模增长趋势',['2020年: 11.76万亿元','2021年: 13.09万亿元 (+11.3%)','2022年: 13.79万亿元 (+5.3%)','2023年: 15.43万亿元 (+11.9%)','2024年: 16.50万亿元 (+6.9%)','2025年: 17.80万亿元 (+7.9% est.)','','数据来源: 国家统计局、商务部',],E)
X(s,0.3,6.7,12.7,0.5,BL)
tt=T(s,0.5,6.75,12.3,0.4);S(tt,'核心洞察: 电商市场持续增长  |  技术门槛是最大障碍  |  AI+鸿蒙带来新机遇  |  中小商家需要低成本方案',fs=13,bo=1,c=P,al=PP_ALIGN.CENTER)

# ===== SLIDE 4: Background 2 - Significance =====
s=bl();H(s,'一、项目背景 — 选题意义')
for i,(x,y,u,v) in enumerate([
    (0.3,1.6,'技术整合实践','SpringBoot+Vue3+鸿蒙全栈集成\nRESTful API + JWT + MyBatis-Plus\nVant移动端 + ElementPlus管理端\n三维设备检测策略(UA+触摸+宽度)'),
    (6.8,1.6,'AI创新落地','DeepSeek大模型电商客服场景应用\n关联词映射引擎(11品类60+关键词)\nRetrieval-Augmented Generation\nAPI失败自动降级(<1s预设回复)'),
    (0.3,4.5,'多端适配工程','移动Web(Vant4 iOS/Android适配)\nPC管理(ElementPlus 响应式适配)\n鸿蒙原生(ArkTS WebView加载)\n京东风格PC端商品详情双栏布局'),
    (6.8,4.5,'商业实用价值','为中小商家提供低门槛解决方案\n完整12张表电商数据模型参考\n40商品/12订单真实初始数据\n阿里云ECS生产环境部署验证'),
]):
    X(s,x,y,6.2,2.6,L)
    X(s,x,y,6.2,0.45,D)
    tt=T(s,x+0.2,y+0.05,5.8,0.35);S(tt,u,fs=16,bo=1,c=W)
    tt=T(s,x+0.2,y+0.6,5.8,1.8)
    lns=v.split('\n');S(tt,lns[0],fs=13,c=D)
    for ln in lns[1:]:P2(tt,ln,fs=12,c=G,sb=6)

# ===== SLIDE 5: Tech Arch Overview =====
s=bl();H(s,'二、技术架构 — 系统全景')
for i,(x,u,v,z) in enumerate([
    (0.3,'前端展示层','Vue3+Vite8构建工具','Vant4(移动端) + ElementPlus(管理端) + ArkTS(鸿蒙)\nPinia状态管理 + VueRouter路由守卫\nAxios HTTP + 请求/响应拦截器\n27个Vue页面 + 7个API模块\nCSS Grid响应式 + SVG占位图'),
    (4.5,'后端服务层','SpringBoot3.2.5 Maven','19个Controller + 11个Service + 12个Mapper\nSpringSecurity6 + JWT(JJWT0.12)无状态认证\nMyBatis-Plus3.5 LambdaWrapper防SQL注入\nDeepSeek RestTemplate AI服务\nApachePOI Excel导出 + Druid监控'),
    (8.7,'基础设施层','MySQL8.0 + CentOS7.9','12张数据表(全部逻辑删除)\nDruid1.2连接池 + SQL监控\n阿里云ECS 2核4G部署\n前后端一体化FatJAR(64MB)\nnohup java -jar 后台运行'),
]):
    X(s,x,1.6,4,5.5,L)
    X(s,x,1.6,4,0.5,D)
    tt=T(s,x+0.2,1.65,3.6,0.4);S(tt,u,fs=16,bo=1,c=W)
    P2(tt,v,fs=11,c=W2,sb=2)
    tt=T(s,x+0.2,2.3,3.6,4.5);lns=z.split('\n');S(tt,lns[0],fs=13,c=D)
    for ln in lns[1:]:P2(tt,ln,fs=12,c=G,sb=5)
X(s,0.3,7,12.7,0.3,L);tt=T(s,0.5,7.02,12.3,0.25)
S(tt,'架构特点: 前后端分离 → RESTful API通信 → JSON数据交互 → JWT无状态认证 → 一体化部署(Fat JAR)',fs=12,bo=1,c=P,al=PP_ALIGN.CENTER)

# ===== SLIDE 6: Frontend Detail =====
s=bl();H(s,'二、技术架构 — 前端技术栈详解')
C(s,0.3,1.5,4,2.8,'移动端用户商城 (Vant4 UI)','组件库: Vant 4.9.24\n布局: Vant NavBar/Tabbar/Sidebar/Tabs\n交互: Swipe轮播/Stepper数量/Popup弹窗\n状态: Pinia cartStore/userStore\n存储: localStorage持久化(购物车/Token)\n适配: 320px-768px触摸优化\n特点: 类原生App体验',P)
C(s,4.5,1.5,4,2.8,'PC管理后台 (ElementPlus)','组件库: Element Plus 2.13.7\n布局: el-container/aside/header/main\n表格: el-table + el-pagination\n表单: el-form + el-dialog CRUD\n图标: @element-plus/icons-vue全集\n移动端: 全局CSS媒体查询响应式\n特点: 深色侧边栏+白色顶栏经典布局',D)
C(s,8.7,1.5,4.3,2.8,'鸿蒙原生端 (ArkTS)','语言: ArkTS 1.2.1 + ArkUI框架\n组件: @kit.ArkWeb Web组件\n功能: 加载云端Web端页面\n配置: domStorageAccess + javaScriptAccess\n适配: HarmonyOS 4.0+\n部署: DevEco Studio构建\n特点: 原生外壳+云端内容',E)
K(s,0.3,4.6,6.2,2.6,'三维设备检测策略 (responsive.js)',['const ua = navigator.userAgent || ""','const isMobileUA = /Android|iPhone|iPad|iPod|Mobile/i.test(ua)','const hasTouch = "ontouchstart" in window || navigator.maxTouchPoints > 0','','function detectMobile() {','  if (isMobileUA) return true  // UA明确移动设备','  if (hasTouch && innerWidth < 1024) return true  // 触摸+窄屏','  if (!hasTouch && innerWidth < 600) return true // 桌面小窗口','  return false','}'],E)
K(s,6.8,4.6,6.2,2.6,'首页并行加载 (Home.vue)',['const loadData = async () => {','  const [categories, hot, new_p, banners] =','    await Promise.all([','      categoryApi.getCategoryList(),','      productApi.getHotProducts(),','      productApi.getNewProducts(),','      bannerApi.list("home").catch(() => [])','    ])','  // 四接口并行请求，提升首屏加载速度','}'],P)

# ===== SLIDE 7: Backend Detail =====
s=bl();H(s,'二、技术架构 — 后端技术栈详解')
for i,(x,u,v,z) in enumerate([
    (0.3,'认证授权模块','Spring Security 6.1 + JWT 0.12','JwtAuthenticationFilter: 拦截Authorization头\n解析Token→提取userId/username/role\n构建UsernamePasswordAuthenticationToken\n@PreAuthorize("hasRole(''ADMIN'')")权限控制\n/error路径permitAll,其余anyRequest.permitAll'),
    (4.5,'持久层模块','MyBatis-Plus 3.5.5 + Druid 1.2','LambdaQueryWrapper: 类型安全查询构建\nBaseMapper: 零XML CRUD操作\n@TableLogic: 全表逻辑删除(deleted=0/1)\n@Select注解: 自定义统计SQL\nDruid监控: /druid/index.html面板'),
    (8.7,'业务与AI模块','11个Service + DeepSeek API','UserService: BCrypt密码加密/资料更新\nProductService: 商品CRUD/热门/新品查询\nOrderInfoService: 订单状态流转/发货\nDeepSeekService: AI对话/商品推荐\nFileUploadController: UUID存储/日期分目录'),
]):
    X(s,x,1.6,4,5.5,L);X(s,x,1.6,4,0.5,D)
    tt=T(s,x+0.2,1.65,3.6,0.4);S(tt,u,fs=16,bo=1,c=W)
    P2(tt,v,fs=11,c=W2,sb=2)
    tt=T(s,x+0.2,2.3,3.6,4.5);lns=z.split('\n');S(tt,lns[0],fs=13,c=D)
    for ln in lns[1:]:P2(tt,ln,fs=12,c=G,sb=5)

# ===== SLIDE 8: Project Stats =====
s=bl();H(s,'二、技术架构 — 项目关键数据')
stats=[(19,'Controller','后端API端点',P),(11,'Service','业务逻辑实现',D),(12,'Mapper','数据访问接口',E),
       (12,'Data Table','数据库表',O),(27,'Vue Page','前端页面',R),(7,'API Module','前端接口封装',P),
       (37,'Test Case','通过率100%',E),('25+','Commit','Git提交记录',D),(3,'Terminal','Web+PC+鸿蒙',O)]
for i,(n,lb,desc,co) in enumerate(stats):
    row=i//3;col=i%3
    SB(s,0.3+col*4.3,1.6+row*2.8,4,2.5,n,lb,co)
    tt=T(s,0.5+col*4.3,1.6+row*2.8+2,3.8,0.3);S(tt,desc,fs=12,bo=1,c=co,al=PP_ALIGN.CENTER)

# ==== SLIDE 9: Requirements - User ====
s=bl();H(s,'三、需求分析 — 用户端功能模块 (10个模块)')
mods=[('🛍 商品浏览','首页Banner+分类导航+搜索\n热门推荐+新品推荐','商品详情','PC京东双栏/移动Vant轮播\n规格选择+价格联动+占位图'),
     ('🛒 购物车','Pinia状态+localStorage\n全选/数量/删除/总金额','📋 订单管理','5状态Tab切换\n取消/支付/确认收货'),
     ('📍 收货地址','CRUD+GPS/IP定位\n设默认+Vant AddressList','👤 个人中心','卡片布局+订单统计\n摄像头拍照换头像'),
     ('⭐ 收藏管理','收藏/取消+列表展示\n点击跳转商品详情','🤖 AI客服','DeepSeek对话\n商品推荐卡片(11品类)'),
     ('🔐 注册登录','JWT+BCrypt+Pinia\n表单校验+路由守卫','💳 下单支付','选地址→确认→提交\n模拟支付弹窗')]
for i,(u,v,x,y) in enumerate(mods):
    row=i//2;col=i%2
    C(s,0.3+col*6.5,1.5+row*2.7,6.2,2.4,u,v,P)
    C(s,3.5+col*6.5,1.5+row*2.7,6.2,2.4,x,y,D)

# ==== SLIDE 10: Requirements - Admin ====
s=bl();H(s,'三、需求分析 — 管理端功能模块 (10个模块)')
for i,(x,v) in enumerate([

    ('📊 仪表盘','四维统计卡片(用户/商品/订单/销售额)\n最近10条订单表格 + 热销Top10列表\nElementPlus Icon彩色图标区'),
    ('👥 用户管理','keyword/role/status筛选搜索 + 分页\nDialog弹窗CRUD(含头像上传)\nBCrypt加密存储 + Excel导出'),
    ('📂 分类管理','el-table展开式一级/二级分类\n名称/排序/状态CRUD\n级联删除确认提示'),
    ('📦 商品管理','三级导航(分类→子分类→商品列表)\n多图上传管理 + keyword/status筛选\n上下架/热销/新品标签 + Excel导出'),
    ('📋 订单管理','orderNo/status筛选 + 分页\nDialog详情弹窗(含order_items明细表)\n待发货→一键发货(status=1→2,记录时间)'),
    ('🎨 Banner管理','CRUD + 图片上传(UUID+日期分目录)\n标题/链接类型/位置/排序/排期设置\nWebMvcConfigurer映射uploads URL路径'),
    ('⭐ 评价管理','productId/status筛选 + 分页\nDialog回复(保存reply+reply_time)\n显示/隐藏切换 + 删除'),
    ('💬 客服消息','用户未读消息列表 + 查看对话记录\n管理员手动回复(sender_type=1)\n未读计数Badge提示'),
    ('⚙ 系统设置','原密码验证 + 新密码BCrypt更新\n系统关于信息(技术栈版本)\n移动端全局CSS响应式适配'),
    ('📱 移动适配','屏幕≤768px: 表格12px/按钮缩小\n弹窗92%宽/抽屉220px\n分页居中/表单项间距缩小'),
]):
    row=i//2;col=i%2
    C(s,0.3+col*6.5,1.5+row*2.5,6.2,2.2,x,v,D)

# ==== SLIDE 11: DB Overview =====
s=bl();H(s,'四、数据库设计 — 12张核心数据表总览')
tables=[('user','15','用户','用户名/BCrypt密码/头像/角色/状态',P),
        ('product','19','商品','名称/价格/库存/销量/热销/新品/详情',P),
        ('category','9','分类','二级结构/base64图标/排序/状态',P),
        ('order_info','16','订单','雪花ID/金额/5状态流转/支付',D),
        ('order_item','11','订单明细','商品快照/单价/数量/小计',D),
        ('product_spec','8','商品规格','SKU(颜色/尺码)/独立价格/库存',E),
        ('address','10','收货地址','省市区/GPS定位/默认标记',E),
        ('banner','11','轮播图','图片/链接/排序/排期/状态',O),
        ('review','13','评价','评分/内容/图片/回复/匿名',O),
        ('favorite','4','收藏','用户-商品关联/时间',R),
        ('cart','7','购物车','用户-商品-规格/数量/选中',R),
        ('customer_service','8','客服','发送者类型/已读标记',G)]
for i,(name,cols,desc,detail,co) in enumerate(tables):
    row=i//4;col=i%4
    X(s,0.3+col*3.2,1.6+row*2.8,3,2.5,L)
    X(s,0.3+col*3.2,1.6+row*2.8,3,0.45,co)
    tt=T(s,0.5+col*3.2,1.65+row*2.8,2.6,0.35);S(tt,f'{name} ({cols}字段)',fs=14,bo=1,c=W)
    tt=T(s,0.5+col*3.2,2.2+row*2.8,2.6,1.5);S(tt,desc,fs=13,bo=1,c=D)
    P2(tt,detail,fs=11,c=G,sb=6)

# ==== SLIDE 12: DB Detail =====
s=bl();H(s,'四、数据库设计 — 设计特点与数据统计')
C(s,0.3,1.5,6.2,2.5,'数据库设计特点','逻辑删除: 全部12张表deleted字段, MyBatis-Plus @TableLogic自动过滤\n订单快照: order_item存储下单时商品名/图片/价格, 防历史数据变更\n雪花ID: Hutool IdUtil.getSnowflakeNextIdStr()生成唯一订单号\nBCrypt: Spring Security BCryptPasswordEncoder加密用户密码\n参数化查询: LambdaQueryWrapper类型安全, 防SQL注入\n字段扩展: 所有表预留扩展字段, 支持灵活业务变更',D)
C(s,6.8,1.5,6.2,2.5,'初始数据统计','商品: 40件 (手机/电脑/家电/服装/美妆/食品/运动等品类)\n订单: 12条 (覆盖5种状态流转)\n订单明细: 18条 (含商品快照)\n评价: 24条 (含评分/图片/管理员回复)\n商品规格: 57条 (颜色/尺码SKU)\n收货地址: 15条 (含默认标记)\nBanner: 5条 (全部配置Unsplash网络图片)\n用户: 8名 (含管理员角色)\n分类: 7个一级分类 + 30个二级分类(含base64图标)',E)
C(s,0.3,4.3,6.2,2.8,'ER关系说明','user 1:N address (收货地址)',D)
P2(T(s,0.5,4.85,5.8,2),'user 1:N order_info (用户订单)\nuser 1:N favorite (收藏)\nuser 1:N review (评价)\nuser 1:N cart (购物车)\nuser 1:N customer_service (客服)',fs=13,c=G,sb=6)
C(s,6.8,4.3,6.2,2.8,'ER关系说明(续)','category 1:N product (商品)',E)
P2(T(s,7,4.85,5.8,2),'product 1:N order_item\nproduct 1:N product_spec\nproduct 1:N review\nproduct 1:N favorite\norder_info 1:N order_item',fs=13,c=G,sb=6)

# ==== SLIDE 13: User Features - Home & Search ====
s=bl();H(s,'五、用户端功能 — 首页与商品浏览','页面路由: /client/home | /client/category | /client/product/:id')
C(s,0.3,1.5,6.2,2.5,'首页功能 (Home.vue)','Banner轮播: VantSwipe autoplay=3000ms, 5条DB数据\n分类导航: 7个base64图标+响应式4-5列网格\n热门商品: /api/products/hot 按销量Top10\n新品推荐: /api/products/new 按时间Top10\n搜索栏: VantSearch圆形搜索框, keyword模糊匹配\n并行加载: Promise.all四接口同时请求',P)
C(s,6.8,1.5,6.2,2.5,'分类与搜索 (Category.vue)','一级分类: VantSidebar侧边栏, 点击切换\n二级分类: VantTabs标签, 动态加载子分类\n商品网格: CSSGrid 3-5列响应式(aspect-ratio:1)\n搜索模式: 隐藏侧边栏, 调用/api/products/search\n占位图: SVG渐变+商品首字(12色轮换)',D)
K(s,0.3,4.3,6.2,2.8,'首页数据加载逻辑',['const [cats, hot, new_p, banners] = await Promise.all([','  categoryApi.getCategoryList(),     // 分类列表','  productApi.getHotProducts(),       // 热门Top10','  productApi.getNewProducts(),       // 新品Top10','  bannerApi.list("home")              // Banner轮播','])','// 四个接口并行请求','// 任一失败不影响其他(individual catch)'],P)
K(s,6.8,4.3,6.2,2.8,'搜索关键词处理',['// 用户输入: "性价比高的手机"','// → 后端提取关键词: ["性价比","手机"]','// → LambdaQueryWrapper:','//   eq(status=1) 上架商品','//   gt(stock>0)   有库存','//   and(or().like(name,keyword)','//          .like(subtitle,keyword)','//          .like(description,keyword))','// → 按销量降序, 分页返回'],D)

# ==== SLIDE 14: User Features - Product Detail ====
s=bl();H(s,'五、用户端功能 — 商品详情(双端自适应)')
C(s,0.3,1.5,6,3,'移动端布局 (Vant4)','顶部: VantSwipe 图片轮播(autoplay=3000)\n价格区: 红色¥价格+原价划线+销量/库存/评分\n服务承诺: 正品保障/7天退换/极速发货\n规格选择: VantCell点击→Popup弹窗\n  └ 标签选择(颜色/尺码) 价格联动切换\n商品详情: HTML富文本v-html渲染\n评价列表: VantRate星级+内容+客服回复\n底部: VantActionBar(客服/收藏/购物车/加购/购买)',P)
C(s,6.5,1.5,6.5,3,'PC端布局 (京东风格)','左侧(450px): 主图(450x450)+缩略图条(60x60)\n  点击缩略图切换主图(border高亮)\n右侧: 商品标题(20px加粗)+副标题(红)\n  粉色价格区块: ¥价格(28px红色加粗)\n  + ¥原价(14px划线)+红色折扣%标签\n  销量/库存/评分 元数据行\n  规格标签选择器 点选切换\n  数量Stepper(min=1 max=库存)\n  大按钮: 加入购物车/立即购买(160x48)\n  + 收藏/客服按钮\n  服务承诺条',D)
K(s,0.3,4.8,6.2,2.4,'规格选择与价格联动',['const selectedSpec = ref(null)','const currentPrice = computed(() =>','  selectedSpec.value?.price || product.price',')','','// 移动端: Popup弹窗 + 标签选择','// PC端: 直接点击标签切换','// 库存联动: max=selectedSpec.stock||product.stock'],P)
K(s,6.8,4.8,6.2,2.4,'PC端缩略图切换',['const images = computed(() => {','  const imgs = []','  if (product.mainImage) imgs.push(...)','  if (product.detailImages) {','    // 兼容JSON数组和逗号分隔','    const arr = JSON.parse(detailImages)','    arr.forEach(img => imgs.push(img))','  }','  return imgs.length ? imgs : [placeholder]','})'],D)

# ==== SLIDE 15: User Features - Cart & Order & Address ====
s=bl();H(s,'五、用户端功能 — 购物车/订单/地址')
C(s,0.3,1.5,4.2,2.5,'购物车 (Cart.vue)\nPinia cartStore','状态: Pinia cartStore集中管理\n存储: localStorage持久化(刷新不丢)\n全选: VantCheckbox全选/取消全选\n数量: VantStepper min=1修改\n删除: 点击删除+确认提示\n结算: VantSubmitBar实时总金额\n跳转: 提交订单→/client/checkout',P)
C(s,4.7,1.5,4.2,2.5,'订单管理 (OrderList.vue)\n5状态流转','Tab: 全部/待付款/待发货/待收货/已完成\n切换: onTabChange()→resetAndLoad()\n  清空列表+重置页码+重新请求\n卡片: 订单号+金额+时间+状态标签\n操作: 查看详情/取消/去支付/确认收货\n详情: /client/order/:id商品明细',D)
C(s,9.1,1.5,4,2.5,'收货地址 (AddressList.vue)\nGPS双重定位','列表: VantAddressList+默认标记\nGPS定位: navigator.geolocation\n  →Nominatim逆地理编码\n  →PERMISSION_DENIED自动回退\n  →ipapi.co→api.ip.sb双重IP定位\n表单: 省/市/区/详细地址+手机号\n操作: CRUD+设置默认',E)
K(s,0.3,4.3,6.2,2.8,'下单流程代码',['async handleSubmit() {','  if (!selectedAddress) { toast("请选地址"); return }','  const items = checkoutItems.map(item => ({','    productId: item.id, quantity: item.quantity','  }))','  const res = await orderApi.createOrder({','    addressId: selectedAddress.id,','    items, remark','  })','  // 返回: { orderId, orderNo, totalAmount }','  // 跳转模拟支付弹窗'],P)
K(s,6.8,4.3,6.2,2.8,'GPS定位回退策略',['async function getLocation() {','  // 1. 尝试GPS','  navigator.geolocation.getCurrentPosition(','    onSuccess,','    onError,  // PERMISSION_DENIED→回退','    { timeout: 8000 }','  )','  // 2. GPS失败→IP定位','  const res = await fetch("ipapi.co/json/")','  // 3. ipapi失败→备用API','  //    api.ip.sb/geoip/'],D)

# ==== SLIDE 16: User Features - Profile & CS ====
s=bl();H(s,'五、用户端功能 — 个人中心/账号设置/AI客服入口')
C(s,0.3,1.5,6.2,2.5,'个人中心 (Profile.vue)\n卡片式布局','用户卡片: 渐变背景+头像+昵称+手机\n订单统计: 5状态数量(API实时):\n  /api/user/orders/counts\n  {all,pendingPay,pendingShip,\n   pendingReceive,completed}\n服务入口: VantGrid宫格\n  我的订单/收货地址/收藏/客服\n操作: 退出登录(确认弹窗)',P)
C(s,6.8,1.5,6.2,2.5,'账号设置 (UserSettings.vue)\n相册+拍照','头像上传: 两个入口\n  1. 相册: <input accept="image/*">\n  2. 拍照: <input capture="camera">\n大小限制: 10MB (FileReader→base64)\n资料编辑: 昵称/手机/邮箱/性别/生日\n性别: ActionSheet选择(男/女)\n生日: DatePicker日期选择器\n保存: /api/user/profile PUT',D)
C(s,0.3,4.3,6.2,2.8,'AI客服入口 (CustomerService.vue)\n每次进入清空历史','欢迎语: "您好，工号26793小琉为您服务"\n快捷问题: 推荐手机/护肤品/新品/包邮/退货\n发送消息: 冒泡动画(三点闪烁)\n  →后端搜索→DeepSeek API→推荐卡片\nAI回复: Markdown符号过滤(正则)\n推荐卡片: 图片+名称+价格+箭头\n  点击跳转/client/product/:id\n未登录: 仍可对话(不保存消息)',E)
K(s,6.8,4.3,6.2,2.8,'客服页核心代码',['const sendMessage = async () => {','  const userMsg = { senderType: 0, content: text }','  messages.value.push(userMsg)','  thinking.value = true','  const res = await csApi.sendMessage(text)','  // 清理Markdown符号','  let reply = res.reply','    .replace(/[*#_`~>-]/g, "")','    .replace(/\\[RECOMMEND\\].*?\\[\\/RECOMMEND\\]/g,"")','  messages.value.push({senderType:1,content:reply,products:res.products})','}'],E)

# ==== SLIDE 17: Admin Features - Dashboard & User ====
s=bl();H(s,'六、管理端功能 — 仪表盘与用户管理')
C(s,0.3,1.5,6.2,3,'仪表盘 (Dashboard.vue)\n/api/admin/dashboard/stats','统计卡片: 4个ElementPlusCard+彩色图标\n  用户总数/商品总数/订单总数/销售总额\n最近订单: el-table展示最近10条\n  (订单号/金额/状态标签)\n热销商品: el-table销量Top10\n  (商品名/销量/价格)\n数据: DashboardStatsVO统一封装\n  {userCount,productCount,orderCount,\n   todayOrderCount,totalSales,todaySales,\n   recentOrders,hotProducts}',P)
C(s,6.8,1.5,6.2,3,'用户管理 (User.vue)\n/api/admin/users','搜索: keyword(用户名/昵称/手机号)+role+status\n列表: ID/用户名/昵称/手机/邮箱/角色标签/状态标签\nCRUD: el-dialog表单(所有字段)\n  用户名/密码/昵称/手机/邮箱\n  头像上传(FileReader→base64)\n  性别/生日/角色/状态\n导出: /api/admin/users/export\n  ApachePOI SXSSFWorkbook\n  自动下载xlsx文件',D)
K(s,0.3,4.8,6.2,2.4,'仪表盘数据查询',['// DashboardServiceImpl.stats()','Long userCount = userMapper.selectCount(wrapper)','Long productCount = productMapper.selectCount(wrapper)','Long orderCount = orderMapper.selectCount(wrapper)','','// 今日订单: count where createTime >= today','Long todayOrders = orderMapper.countTodayOrders(start,end)','// 销售额: select ifnull(sum(pay_amount),0)','BigDecimal totalSales = orderMapper.sumPayAmount()'],P)
K(s,6.8,4.8,6.2,2.4,'用户导出Excel',['// UserAdminController.export()','@GetMapping("/export")','public void export(HttpServletResponse response) {','  List<User> users = userService.listAll(...)','  String[] headers = {"ID","用户名","昵称",...}','  List<List<Object>> data = users.stream()','    .map(u -> Arrays.asList(u.getId(),...))','    .collect(toList())','  ExcelExportUtil.exportExcel(response,headers,data)','}'],D)

# ==== SLIDE 18: Admin Features - Category & Product ====
s=bl();H(s,'六、管理端功能 — 分类管理与商品管理')
C(s,0.3,1.5,6.2,2.8,'分类管理 (Category.vue)\n/api/admin/categories','一级分类: el-table + type="expand"\n  展开显示二级子分类表格\n  ID/名称/排序/状态标签/操作\n二级分类: expand-content区域\n  独立子表格(名称/排序/状态)\nDialog: 新增/编辑(名称/排序/状态)\n  名称必填(trigger:blur校验)\n删除: ElMessageBox.confirm确认\n  级联删除提示(子分类+商品)',P)
C(s,6.8,1.5,6.2,2.8,'商品管理 (Product.vue)\n/api/admin/products','三级导航: 一级分类→子分类→商品列表\n  (el-breadcrumb面包屑导航)\n列表: keyword/status筛选+分页\n  主图缩略图/名称/价格/库存/销量\nDialog: 18个字段完整表单\n  名称/副标题/主图/详情图片(多张)\n  售价/原价/库存/销量/排序\n  描述/详情HTML/状态/热销/新品\n导出: /api/admin/products/export',D)
K(s,0.3,4.6,6.2,2.6,'商品数据查询',['// ProductAdminController.page()','LambdaQueryWrapper<Product> wrapper','  = new LambdaQueryWrapper<>()','wrapper.eq(categoryId!=null, Product::getCategoryId, id)','  .like(StrUtil.isNotBlank(keyword),','         Product::getName, keyword)','  .eq(status!=null, Product::getStatus, status)','  .orderByDesc(Product::getSortOrder)','Page<Product> page = mapper.selectPage(',  '  new Page<>(pageNum,pageSize), wrapper)'],P)
K(s,6.8,4.6,6.2,2.6,'分类展开子分类',['const getChildren = (parentId) => {','  if (!parentId) return []','  if (!childrenCache.value[parentId]) {','    // 首次展开→加载数据','    request.get(`/api/admin/categories/', '      second/${parentId}`).then(res => {','      // 对象替换触发Vue响应式更新','      childrenCache.value = {','        ...childrenCache.value,','        [parentId]: res || []','      }','    })','}'],D)

# ==== SLIDE 19: Admin Features - Order & Banner ====
s=bl();H(s,'六、管理端功能 — 订单管理与Banner管理')
C(s,0.3,1.5,6.2,2.8,'订单管理 (Order.vue)\n/api/admin/orders','搜索: orderNo/status筛选+分页\n列表: ID/订单号/金额/支付方式/状态标签/时间\n操作: 查看(详情弹窗)/发货/导出\nDialog: el-descriptions订单信息\n  (订单号/金额/实付/支付方式/时间)\n  商品明细table: 商品名/规格/单价/数量/小计\n发货: /api/admin/orders/{id}/ship\n  status=1→2(待收货), record delivery_time\n导出: /api/admin/orders/export→xlsx',P)
C(s,6.8,1.5,6.2,2.8,'Banner管理 (Banner.vue)\n/api/admin/banners','列表: 图片缩略图/标题/位置/排序/状态/时间\nDialog: 标题(input)/图片(upload预览)\n  链接类型(商品/分类/URL select)\n  链接目标(input)/位置(home)\n  排序(number)/开始时间/结束时间(datepicker)\n  状态(启用/禁用 select)\n上传: uploadApi.uploadImage(file)\n  →FormData→/api/upload/image\n  →UUID命名→日期分目录→返回URL',D)
K(s,0.3,4.6,6.2,2.6,'订单状态流转逻辑',['@PutMapping("/{id}/ship")  // 发货','public ApiResponse<Void> ship(@PathVariable Long id) {','  OrderInfo order = orderInfoService.getById(id)','  if (order != null && order.getStatus() == 1) {','    order.setStatus(2)  // 待收货','    order.setDeliveryTime(LocalDateTime.now())','    orderInfoService.update(order)','  }','  return ApiResponse.success()','}'],P)
K(s,6.8,4.6,6.2,2.6,'图片上传处理',['@PostMapping("/image")','public ApiResponse uploadImage(@RequestParam file){','  if (!contentType.startsWith("image/"))','    return error("只支持图片文件")','  String dateDir = now().format("yyyy/MM/dd")','  File datePath = new File(uploadDir, dateDir)','  if (!datePath.exists()) datePath.mkdirs()','  String name = UUID.randomUUID()+ext','  file.transferTo(new File(datePath, name))','  return success("/uploads/"+dateDir+"/"+name)','}'],D)

# ==== SLIDE 20: Admin Features - Review & Settings ====
s=bl();H(s,'六、管理端功能 — 评价管理与系统设置')
C(s,0.3,1.5,6.2,3,'评价管理 (Review.vue)\n/api/admin/reviews','搜索: productId+status筛选+分页\n列表: ID/商品ID/用户ID/VantRate星级/内容\n  回复内容/状态标签/时间\nDialog: 回复(展示原评价+textarea输入)\n  保存: reply+reply_time字段\n状态: 显示/隐藏切换(updateStatus)\n删除: 确认弹窗+删除+刷新列表',P)
C(s,6.8,1.5,6.2,3,'系统设置 (Settings.vue)\n/api/user/profile','修改密码: el-tabs切换\n  原密码/新密码/确认密码\n  校验: 新密码≥6位,两次一致\n  →/api/user/profile PUT\n  →BCrypt加密更新\n  →清除Token重新登录\n关于: 系统版本号+技术栈信息\n客服消息: 用户对话查看+回复\n移动适配: 全站CSS媒体查询\n  ≤768px响应式规则',D)
K(s,0.3,4.8,6.2,2.4,'评价回复处理',['@PutMapping("/{id}/reply")','public ApiResponse reply(','    @PathVariable Long id,','    @RequestBody Map<String,String> body) {','  reviewService.reply(id, body.get("reply"))','','// ReviewServiceImpl.reply()','Review review = new Review()','review.setId(id)','review.setReply(reply)','review.setReplyTime(LocalDateTime.now())','reviewMapper.updateById(review)'],P)
K(s,6.8,4.8,6.2,2.4,'移动端CSS响应式',['/* App.vue 全局媒体查询 */','@media (max-width: 768px) {','  .el-table { font-size: 12px !important }','  .el-table .cell { padding: 6px 4px }','  .el-button--small {','    padding: 4px 8px; font-size: 11px }','  .el-form--inline .el-form-item {','    margin-right: 8px; margin-bottom: 4px }','  .el-pagination { justify-content: center }','  .el-dialog { width: 92%!important }','  .el-drawer { width: 220px!important }}'],D)

# ==== SLIDE 21: AI Innovation Detail =====
s=bl();H(s,'七、AI智能客服 — 核心创新点详解 ⭐')
C(s,0.3,1.5,6.2,2.5,'技术方案\nDeepSeek+RAG简化实现','模型: deepseek-chat (API via RestTemplate)\nSystemPrompt: 动态构建(含商品列表+规则)\n  禁止Markdown符号,口语化要求\n  角色: "QJ商城AI助手小Q,工号26793"\n关联词映射: HashMap<String,String[]>\n  11个品类 60+关键词\n  饮料→气泡水/可乐/橙汁/果汁...\n  手机→iPhone/华为/小米/OPPO...',E)
C(s,6.8,1.5,6.2,2.5,'工作流程\n6步处理链路','Step1: 保存用户消息→提取关键词\nStep2: 扩展关联词→搜索DB商品(最多10)\nStep3: 取最近10条对话历史\nStep4: 构建SystemPrompt+history+userMsg\nStep5: 调用DeepSeek API\n        headers.setBearerAuth(apiKey)\n        model=deepseek-chat,temperature=0.7\nStep6: 解析[RECOMMEND]标记→组装卡片\n        失败→catch→fallback预设回复(<1s)',O)
K(s,0.3,4.3,6.2,2.8,'关联词映射引擎',['static final Map<String,String[]> RELATED = {','  {"饮料",{"气泡水","可乐","橙汁","果汁"}},','  {"手机",{"iPhone","华为","小米","三星"}},','  {"电脑",{"MacBook","ThinkPad","笔记本"}},','  {"耳机",{"AirPods","Bose","索尼","降噪"}},','  {"零食",{"坚果","车厘子","食品"}},','  {"衣服",{"羽绒","卫衣","CK","优衣库","Nike"}},','  {"护肤",{"SK-II","雅诗兰黛","海蓝之谜"}},','  // 共11品类60+关键词'],E)
K(s,6.8,4.3,6.2,2.8,'推荐解析与降级',['// 解析AI回复中的推荐标记','if (content.contains("[RECOMMEND]")) {','  int s = content.indexOf("[RECOMMEND]")+11','  int e = content.indexOf("[/RECOMMEND]")','  String ids = content.substring(s,e)','  // "108,106" → [108L, 106L]','  // 最多3个','  // 组装productCards: {id,name,price,image,link}','','// API失败→降级','} catch (Exception e) {','  return fallbackReply(userMessage)  // 预设回复','}'],O)

# ==== SLIDE 22: AI Innovation Detail 2 =====
s=bl();H(s,'七、AI智能客服 — 前端展示与效果示例')
C(s,0.3,1.5,6.2,2.8,'前端聊天UI (CustomerService.vue)','欢迎区: 大号机器人emoji + 工号 + 欢迎语\n  5个快捷问题标签(推荐手机/护肤品等)\n消息气泡: 用户(蓝色右对齐) / AI(白色左对齐)\n  AI思考中: 三点闪烁动画(blink keyframes)\n  msg.content: Markdown符号正则过滤\n  msg.products: 商品推荐卡片数组\n    └ 图片(56x44圆角)+名称+价格+箭头\n    └ @click="goProduct(p.id)"跳转详情\n输入栏: VantField+发送按钮(圆形)',P)
C(s,6.8,1.5,6.2,2.8,'效果演示','用户: "我想要饮料"',D)
P2(T(s,7.2,2.2,5.4,1),'AI: "为您推荐以下饮料哦~"',
   fs=13,c=E,sb=2)
P2(T(s,7.2,2.4,5.4,1),'【元气森林气泡水】无糖低卡 480ml*15 \n    ¥59.90  [商品卡片可点击]',
   fs=12,c=P,sb=2)
P2(T(s,7.2,2.9,5.4,1),'【智利进口车厘子】JJ级果径饱满 2.5kg \n    ¥199.00  [商品卡片可点击]',
   fs=12,c=P,sb=2)
P2(T(s,7.2,3.4,5.4,1),'用户: "推荐一款手机"',
   fs=13,c=D,sb=6)
P2(T(s,7.2,3.6,5.4,1),'AI: "为您推荐以下手机~"',
   fs=13,c=E,sb=2)
P2(T(s,7.2,3.8,5.4,1),'【iPhone 15 Pro Max】A17 Pro芯片 ¥9999\n【小米14 Pro】徕卡影像 ¥4999\n【三星S24 Ultra】AI手机 ¥9699',
   fs=12,c=P,sb=2)
K(s,0.3,4.6,6.2,2.6,'前端消息处理',['const sendMessage = async () => {','  const text = inputText.value.trim()','  if (!text || thinking.value) return','  thinking.value = true','  messages.push({content:text,senderType:0})','  const res = await csApi.sendMessage(text)','  let reply = res.reply','    .replace(/[*#_`~>-]/g, "")  // 去符号','    .replace(/\\[RECOMMEND\\].*?\\[\\/RECOMMEND\\]/g,"")','  messages.push({content:reply,senderType:1,products:res.products})','}'],P)
K(s,6.8,4.6,6.2,2.6,'后端商品搜索',['private List<Product> searchProducts(String keyword) {','  Set<String> words = extractAndExpand(keyword)','  LambdaQueryWrapper<Product> w = new LambdaQueryWrapper<>()','  w.eq(Product::getStatus, 1)  // 上架','  w.gt(Product::getStock, 0)   // 有库存','  w.and(wr -> {','    for (String word : words)','      wr.or().like(Product::getName, word)','        .or().like(Product::getSubtitle, word)','        .or().like(Product::getDescription, word)','  })','  w.orderByDesc(Product::getSales)','  w.last("limit 10")','}'],D)

# ==== SLIDE 23: Testing Overview =====
s=bl();H(s,'八、作品测试 — 测试概览')
SB(s,0.3,1.5,3.8,1.2,'37','测试用例总数',P)
SB(s,4.3,1.5,3.8,1.2,'100%','测试通过率',E)
SB(s,8.3,1.5,3.8,1.2,'2类','用户端+管理端',D)
tt=T(s,0.3,3,12.7,0.5);S(tt,'测试环境',fs=18,bo=1,c=D)
P2(tt,'服务器: 阿里云ECS 2核4G CentOS 7.9 | 数据库: MySQL 8.0 (qj_shop) | 浏览器: Chrome 120+ / Edge 120+ / Safari 17+\n移动设备: iPhone 14 Pro (390×844) / Android Chrome (412×915) | PC: 1920×1080 / 1366×768 | 鸿蒙: DevEco模拟器 HarmonyOS 4.0\n初始数据: 40件商品 | 12条订单 | 24条评价 | 8名用户 | 7个一级分类 | 5条Banner',fs=13,c=G,sb=6)
C(s,0.3,3.8,6.2,3.2,'用户端测试用例(24个)','注册(合法/重复) ✓  登录(正确/错误) ✓\n首页加载(Banner/分类/热门) ✓  搜索(关键词模糊) ✓\n商品详情(移动/PC双端) ✓  规格选择(价格联动) ✓\n购物车(加入/修改/删除/全选) ✓\n下单(选地址→确认→提交→支付) ✓\n订单(5状态Tab切换/取消/支付) ✓\n收货地址(GPS定位/编辑/默认) ✓\nAI客服(泛指推荐/口语回复/降级) ✓\n个人中心(统计/头像拍照/资料) ✓',E)
C(s,6.8,3.8,6.2,3.2,'管理端测试用例(13个)','仪表盘(统计/最近订单/热销) ✓\n用户管理(CRUD/搜索/Excel导出) ✓\n分类管理(展开/CRUD/级联删除) ✓\n商品管理(三级/筛选/上下架/导出) ✓\n订单管理(筛选/详情弹窗/发货) ✓\nBanner管理(CRUD/图片上传) ✓\n评价管理(回复/显示隐藏/删除) ✓\n移动端管理(表格/弹窗/分页响应) ✓',P)

# ==== SLIDE 24: Testing Details =====
s=bl();H(s,'八、作品测试 — 兼容性与性能测试')
C(s,0.3,1.5,6.2,2.8,'浏览器兼容性','Chrome 120+: 完整通过 ✓\nEdge 120+: 完整通过 ✓\nSafari 17+: 完整通过 ✓\n微信内置浏览器(Android): 通过 ✓\n微信内置浏览器(iOS): 通过 ✓\n鸿蒙DevEco模拟器: WebView正常加载 ✓\n响应式测试: 320px~1920px全覆盖 ✓\n  小屏(<500px): 商品3列网格\n  中屏(500-900px): 4列\n  大屏(>900px): 5列',P)
C(s,6.8,1.5,6.2,2.8,'性能指标','页面首次加载: <3秒 ✓\n商品列表查询: <1秒 ✓\n订单提交: <2秒 ✓\nAI客服(正常): 3-5秒 ✓\nAI客服(降级): <1秒 ✓\n并发50用户浏览: 无明显延迟 ✓\nJAR启动时间: ~7秒\nJAR大小: 64MB\n前端dist: 38文件\n后端class: 80+文件',D)
C(s,0.3,4.6,6.2,2.6,'设备检测准确性验证','UA检测: /Android|iPhone|iPad|Mobile/i\n  → 正确识别100%移动UA ✓\n触摸屏检测: ontouchstart + maxTouchPoints\n  → 正确识别触摸设备 ✓\n宽度兜底: <600px强制移动\n  → 窄窗口桌面浏览器正确切换 ✓\n高分辨率手机: 1080px Android旗舰\n  → UA→移动端, 不再误判PC ✓\niPad: 触摸屏+中等宽度\n  → 识别为平板, 使用移动布局 ✓',E)
K(s,6.8,4.6,6.2,2.6,'性能监控代码',['// Druid连接池监控','// /druid/index.html','// 实时查看:','//   - 活跃连接数','//   - SQL执行时间','//   - 慢SQL记录(>2s)','// ','// MyBatis-Plus SQL日志','// mybatis-plus.configuration','//   .log-impl=StdOutImpl','// ','// 每个API请求打印SQL和参数','// 便于性能分析和优化'],D)

# ==== SLIDE 25: Problems 1 ====
s=bl();H(s,'九、技术问题与解决 (Page 1/2)')
problems1=[('问题1','Vant4组件名变更','GoodsAction系列→ActionBar\nProductDetail组件解析失败','查阅Vant官方CHANGELOG\nvan-goods-action→van-action-bar\nvan-goods-action-icon→van-action-bar-icon\nvan-goods-action-button→van-action-bar-button'),
('问题2','高分辨率手机误判PC','1080px Android手机\n仅靠innerWidth显示PC布局','UA检测+触摸屏检测+宽度三维策略\nisMobileUA || (hasTouch && w<1024)'),
('问题3','Vite代理rewrite路径','代理去/api→Controller\n@RequestMapping("/api/...")不匹配','移除rewrite规则\n去掉axios baseURL\n所有API路径统一完整/api/...'),
('问题4','Banner时间过滤NULL','startTime IS NULL\n→ startTime<=now=NULL→过滤','SQL增加IS NULL OR判断\nisNull(startTime).or().le(startTime,now)')]
for i,(no,name,prob,sol) in enumerate(problems1):
    row=i//2;col=i%2
    X(s,0.3+col*6.5,1.5+row*2.8,6.2,2.5,L)
    X(s,0.3+col*6.5,1.5+row*2.8,6.2,0.45,R)
    tt=T(s,0.5+col*6.5,1.55+row*2.8,5.8,0.35);S(tt,no+' - '+name,fs=14,bo=1,c=W)
    tt=T(s,0.5+col*6.5,2.1+row*2.8,5.8,0.8);S(tt,'现象: '+prob,fs=12,c=D)
    P2(tt,'解决: '+sol,fs=12,c=E,sb=6)

# ==== SLIDE 26: Problems 2 ====
s=bl();H(s,'九、技术问题与解决 (Page 2/2)')
problems2=[('问题5','Linux MySQL Socket','JDBC:localhost→UnixSocket\nsocket文件不在默认路径','改为jdbc:mysql://127.0.0.1:3306/\n强制TCP连接'),
('问题6','SPA路由刷新500','/admin直接访问→404\n→500服务器异常','SpaFilter@Order(MIN_VALUE)\n非API请求forward:/index.html\n保留query参数'),
('问题7','文件上传Content-Type','手动设multipart/form-data\n缺boundary, 服务端无法解析','去除手动Content-Type设置\n让axios自动生成含boundary的请求头'),
('问题8','注册密码双重BCrypt','Controller+Service各自encode\n存储BCrypt(BCrypt(明文))','去掉Controller重复编码\n统一由UserServiceImpl.save()处理')]
for i,(no,name,prob,sol) in enumerate(problems2):
    row=i//2;col=i%2
    X(s,0.3+col*6.5,1.5+row*2.8,6.2,2.5,L)
    X(s,0.3+col*6.5,1.5+row*2.8,6.2,0.45,O)
    tt=T(s,0.5+col*6.5,1.55+row*2.8,5.8,0.35);S(tt,no+' - '+name,fs=14,bo=1,c=W)
    tt=T(s,0.5+col*6.5,2.1+row*2.8,5.8,1);S(tt,'现象: '+prob,fs=12,c=D)
    P2(tt,'解决: '+sol,fs=12,c=E,sb=6)
X(s,0.3,5.8,12.7,1.2,BL);tt=T(s,0.5,5.9,12.3,1)
S(tt,'经验总结',fs=16,bo=1,c=W)
P2(tt,'• 文档先行: 升级依赖前查阅CHANGELOG和Migration Guide    • 编码规范: 统一加密/验证在Service层,Controller只做参数接收\n• 防御编程: 数据库查询考虑NULL边界    • 代理配置: 理解前端代理与后端路径的映射关系    • 生产部署: 优先TCP连接避免Socket路径问题',fs=13,c=W2,sb=8)

# ==== SLIDE 27: Summary - Achievements ====
s=bl();H(s,'十、项目总结 — 成果汇总')
achievements=[('80+文件','源代码总数','涵盖后端Java+前端Vue+鸿蒙ArkTS+配置+SQL+文档',P),
              ('19 Controller','后端API端点','5个公共+5个客户端+9个管理端, 覆盖全部业务场景',D),
              ('12张数据表','MySQL 8.0','从用户/商品/订单核心表到地址/收藏/评价扩展表',E),
              ('27个Vue页面','前端组件','15个客户端页面+12个管理端页面, 双端独立',O),
              ('3个终端','多端覆盖','移动Web(Vant)+PC管理(ElementPlus)+鸿蒙(ArkTS)',R),
              ('37个测试','100%通过','用户端24个+管理端13个, 覆盖全部核心功能',P)]
for i,(n,lb,desc,co) in enumerate(achievements):
    X(s,0.3,1.5+i*0.95,12.7,0.85,L)
    X(s,0.3,1.5+i*0.95,1.5,0.85,co)
    tt=T(s,0.4,1.55+i*0.95,1.3,0.35);S(tt,n,fs=16,bo=1,c=W,al=PP_ALIGN.CENTER)
    tt=T(s,2.0,1.55+i*0.95,10.8,0.35);S(tt,lb,fs=14,bo=1,c=D)
    P2(tt,desc,fs=12,c=G,sb=2)

# ==== SLIDE 28: Summary - Innovation ====
s=bl();H(s,'十、项目总结 — 创新点')
for i,(u,v,z) in enumerate([
    ('AI客服 + RAG简化实现','DeepSeek大模型 + 关联词映射引擎(11品类60+关键词)\n自然语言→商品推荐(最多3个)→聊天卡片→点击购买\nAPI失败自动降级(<1s预设回复)','技术难度: ★★★★☆ | 实用价值: ★★★★★'),
    ('三维设备检测策略','UA检测 + 触摸屏检测 + 屏幕宽度三维判定\n解决高分辨率手机(1080px)误判PC端问题\n移动端Vant4+PC端京东风格, 自动切换无需用户干预','技术难度: ★★★☆☆ | 实用价值: ★★★★☆'),
    ('GPS/IP双重定位','GPS(W3CGeolocation)→PermissionDenied→IP\nNominatim逆地理编码(含User-Agent头)\n→ipapi.co→api.ip.sb双重IP回退','技术难度: ★★★☆☆ | 实用价值: ★★★★☆'),
    ('前后端一体化部署','Vue前端dist→SpringBoot static目录→Maven Package\nFatJAR 64MB, java -jar一键启动\n无需Nginx/分离部署, 降低运维复杂度','技术难度: ★★☆☆☆ | 实用价值: ★★★★★'),
    ('12张表完整数据模型','覆盖user→order_info→order_item完整订单链路\n快照机制(order_item存储下单时商品快照)\n逻辑删除(@TableLogic)保证数据安全可恢复','技术难度: ★★★☆☆ | 实用价值: ★★★★★'),
]):
    row=i//2;col=i%2
    X(s,0.3+col*6.5,1.5+row*1.8,6.2,1.6,L)
    X(s,0.3+col*6.5,1.5+row*1.8,6.2,0.42,A)
    tt=T(s,0.5+col*6.5,1.55+row*1.8,5.8,0.35);S(tt,u,fs=14,bo=1,c=W)
    tt=T(s,0.5+col*6.5,2.05+row*1.8,5.8,0.8);S(tt,v.split('\n')[0],fs=12,c=D)
    for ln in v.split('\n')[1:]:P2(tt,ln,fs=11,c=G,sb=3)
    P2(tt,z,fs=10,c=O,sb=4)

# ==== SLIDE 29: Summary - Future ====
s=bl();H(s,'十、项目总结 — 不足与改进方向')
for i,(x,v,c) in enumerate([

    ('支付真实化','当前: 模拟支付按钮+状态流转(前端确认)\n改进: 接入支付宝/微信支付SDK\n  → 真实支付回调+订单状态自动更新\n  → 支付安全(签名验证/金额校验)',P),
    ('AI推荐精准化','当前: 关联词映射(品类→关键词)搜索DB\n改进: 引入向量数据库(Milvus/Weaviate)\n  → 商品描述embedding语义搜索\n  → 用户行为推荐(协同过滤/深度学习)',D),
    ('图片存储优化','当前: base64存储在数据库longtext字段\n  → 数据库体积大, 查询效率低\n改进: 迁移至OSS(阿里云/七牛)对象存储\n  → 图片URL引用, 按需加载',E),
    ('性能优化','当前: 单体应用, 单机部署(2核4G)\n改进: Redis缓存热门数据(商品/分类)\n  → 消息队列(RabbitMQ)异步处理订单\n  → 分布式架构(SpringCloud微服务化)',O),
    ('功能扩展','当前: 覆盖电商核心流程\n改进: 优惠券/满减/秒杀活动\n  → 物流追踪(快递100 API)\n  → 数据可视化(ECharts销售趋势图)\n  → 消息推送(WebSocket实时通知)',R),
    ('安全增强','当前: BCrypt+JWT+@PreAuthorize\n改进: HTTPS/TLS证书配置\n  → 验证码(滑块/短信)防刷\n  → 接口限流(RateLimiter/令牌桶)\n  → SQL审计日志+异常IP检测',G),
]):
    row=i//3;col=i%3
    X(s,0.3+col*4.3,1.5+row*2.8,4,2.5,L)
    X(s,0.3+col*4.3,1.5+row*2.8,4,0.45,c)
    tt=T(s,0.5+col*4.3,1.55+row*2.8,3.6,0.35);S(tt,x,fs=14,bo=1,c=W)
    tt=T(s,0.5+col*4.3,2.05+row*2.8,3.6,1.7);lns=v.split('\n');S(tt,lns[0],fs=12,c=D)
    for ln in lns[1:]:P2(tt,ln,fs=11,c=G,sb=5)

# ==== SLIDE 30: Thanks =====
s=bl();X(s,0,0,13.333,7.5,P);X(s,0,0,13.333,0.08,A);X(s,0,7.42,13.333,0.08,A)
tt=T(s,1.5,1.5,10.3,1.5);S(tt,'感谢聆听',fs=56,bo=1,c=W,al=PP_ALIGN.CENTER)
P2(tt,'请各位老师批评指正',fs=26,c=W2,al=PP_ALIGN.CENTER,sb=16)
X(s,4,3.5,5.3,0.01,W)
tt=T(s,3,4,7.3,3)
S(tt,'"柒玖商店" — 基于SpringBoot多前端多平台的商城系统',fs=16,c=W,al=PP_ALIGN.CENTER)
for ln,sz in [('',10),('姓   名：李  鑫',16),('学   号：202330302152',16),('专   业：软件技术  |  班 级：2307',16),('指导教师：彭德宇',16),('2026年5月',14)]:
    P2(tt,ln,fs=sz,c=W if sz>14 else W2,al=PP_ALIGN.CENTER,sb=8)

# ===== SAVE =====
path=r"D:\QJShop\柒玖商店_毕业设计答辩_final.pptx"
prs.save(path)
print(f"SAVED: {path}")
print(f"SLIDES: {len(prs.slides)}")
