from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

# ===== Color Palette =====
C_PRIMARY   = RGBColor(0x25, 0x63, 0xEB)  # Blue
C_DARK      = RGBColor(0x1A, 0x36, 0x5D)  # Navy
C_ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)  # Gold
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x00, 0x00, 0x00)
C_GRAY      = RGBColor(0x64, 0x74, 0x8B)  # Slate
C_LIGHT     = RGBColor(0xF1, 0xF5, 0xF9)  # Light bg
C_RED       = RGBColor(0xEF, 0x44, 0x44)
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)
C_BLUE_LT   = RGBColor(0xDB, 0xEA, 0xFE)  # Light blue
C_ORANGE    = RGBColor(0xF9, 0x73, 0x16)

W = prs.slide_width
H = prs.slide_height

# ===== Helper Functions =====
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame

def set_text(tf, text, font_name='Microsoft YaHei', font_size=18, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

def add_para(tf, text, font_name='Microsoft YaHei', font_size=14, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_rect(slide, left, top, width, height, color, radius=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ===== Slide 1: Cover =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 7.5, C_PRIMARY)
add_rect(s, 0, 0, 13.333, 0.08, C_ACCENT)
add_rect(s, 0, 7.42, 13.333, 0.08, C_ACCENT)

tb = add_textbox(s, 1.5, 1.0, 10.3, 1.8)
set_text(tb, '"柒玖商店"', font_size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_para(tb, '基于SpringBoot多前端多平台的商城系统', font_size=28, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER, space_before=12)

add_rect(s, 4, 3.2, 5.3, 0.01, C_WHITE)

tb = add_textbox(s, 4.5, 3.6, 4.5, 2.8)
set_text(tb, '毕业设计答辩', font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_para(tb, '姓   名：李  鑫', font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=16)
add_para(tb, '学   号：202330302152', font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=8)
add_para(tb, '专   业：软件技术  |  班 级：2307', font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=8)
add_para(tb, '指导教师：彭德宇', font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=8)
add_para(tb, '2026年5月', font_size=16, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER, space_before=20)

# ===== Slide 2: Table of Contents =====
s = add_blank_slide()
add_rect(s, 0, 0, 0.08, 7.5, C_PRIMARY)
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
tb = add_textbox(s, 0.8, 0.4, 11, 0.8)
set_text(tb, '答辩大纲', font_size=36, bold=True, color=C_DARK)

items = [
    ('01', '项目背景与选题意义', '电子商务发展趋势与课题价值'),
    ('02', '技术架构概述', '前后端技术选型与系统架构'),
    ('03', '需求分析与功能模块', '用户端+管理端+鸿蒙端功能矩阵'),
    ('04', '数据库设计', '12张核心数据表与ER关系'),
    ('05', '核心功能演示 - 用户端', '商品浏览/下单/AI客服等'),
    ('06', '核心功能演示 - 管理端', '仪表盘/商品管理/订单管理等'),
    ('07', 'AI智能客服创新点', 'DeepSeek集成+RAG实现+关联推荐'),
    ('08', '作品测试', '37个测试用例+兼容性验证'),
    ('09', '项目总结与展望', '成果/创新/不足/改进方向'),
]
y = 1.6
for num, title, desc in items:
    add_circle(s, 0.7, y + 0.12, 0.45, C_PRIMARY)
    ct = add_textbox(s, 1.02, y + 0.15, 0.35, 0.4)
    set_text(ct, num, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    tt = add_textbox(s, 1.5, y, 7, 0.35)
    set_text(tt, title, font_size=16, bold=True, color=C_DARK)
    dt = add_textbox(s, 1.5, y + 0.32, 9, 0.25)
    set_text(dt, desc, font_size=11, color=C_GRAY)
    y += 0.62

# ===== Slide 3: Background =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '一、项目背景与选题意义', font_size=32, bold=True, color=C_WHITE)

add_rect(s, 0.5, 1.6, 5.8, 0.5, C_BLUE_LT)
tb = add_textbox(s, 0.7, 1.65, 5.4, 0.4)
set_text(tb, '📊 选题背景', font_size=20, bold=True, color=C_DARK)

tb = add_textbox(s, 0.7, 2.3, 5.3, 4)
set_text(tb, '• 电子商务持续增长，2025年中国网络零售市场超15万亿元', font_size=15, color=C_DARK)
add_para(tb, '• 中小型商家数字化转型门槛高、成本大，自建平台技术难度大', font_size=15, color=C_DARK, space_before=12)
add_para(tb, '• AI大模型(DeepSeek等)快速崛起，为电商智能化带来新可能', font_size=15, color=C_DARK, space_before=12)
add_para(tb, '• 鸿蒙生态逐步完善，应用需要多终端覆盖能力', font_size=15, color=C_DARK, space_before=12)
add_para(tb, '• 亟需低成本、易维护的中小型电商解决方案', font_size=15, color=C_DARK, space_before=12)

add_rect(s, 7.0, 1.6, 5.8, 0.5, C_BLUE_LT)
tb = add_textbox(s, 7.2, 1.65, 5.4, 0.4)
set_text(tb, '🎯 选题意义', font_size=20, bold=True, color=C_DARK)

tb = add_textbox(s, 7.2, 2.3, 5.3, 4)
set_text(tb, '✦ 技术整合：Spring Boot+Vue3+鸿蒙的全栈工程实践', font_size=15, color=C_DARK)
add_para(tb, '✦ AI落地：DeepSeek大模型在电商客服场景的真实应用', font_size=15, color=C_DARK, space_before=12)
add_para(tb, '✦ 多端适配：移动Web+PC管理+鸿蒙三端覆盖', font_size=15, color=C_DARK, space_before=12)
add_para(tb, '✦ 实用价值：为中小商家提供低成本数字化转型参考', font_size=15, color=C_DARK, space_before=12)

add_rect(s, 0.5, 5.8, 12.3, 1.2, RGBColor(0xFE,0xFB,0xEB))
add_circle(s, 0.8, 6.1, 0.4, C_ACCENT)
tb = add_textbox(s, 1.3, 6.05, 11, 0.9)
set_text(tb, '核心亮点', font_size=16, bold=True, color=C_ORANGE)
add_para(tb, 'AI客服(11品类60+关键词联想推荐)  |  双端自适应(京东风格PC+移动Vant)  |  GPS/IP双重定位  |  12张表完整数据模型', font_size=13, color=C_DARK, space_before=6)

# ===== Slide 4: Tech Architecture =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '二、技术架构概述', font_size=32, bold=True, color=C_WHITE)

# Frontend
add_rect(s, 0.3, 1.5, 4.0, 0.45, C_DARK)
tb = add_textbox(s, 0.5, 1.53, 3.6, 0.4)
set_text(tb, '🖥 前端技术栈', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 0.5, 2.1, 3.8, 4.5)
set_text(tb, '移动端用户商城:', font_size=15, bold=True, color=C_DARK)
add_para(tb, 'Vue 3.5 + Vite 8 + Vant 4 UI', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, 'Pinia状态管理 | Axios HTTP', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, 'PC端管理后台:', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'Vue 3.5 + Element Plus 2.13', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, '全局CSS变量 | 响应式适配', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, '设备检测:', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'UA + 触摸屏 + 屏幕宽度\n三维判定策略', font_size=13, color=C_GRAY, space_before=6)

# Backend
add_rect(s, 4.6, 1.5, 4.0, 0.45, C_DARK)
tb = add_textbox(s, 4.8, 1.53, 3.6, 0.4)
set_text(tb, '⚙ 后端技术栈', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 4.8, 2.1, 3.8, 4.5)
set_text(tb, '核心框架:', font_size=15, bold=True, color=C_DARK)
add_para(tb, 'Spring Boot 3.2.5 + Maven', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, 'MyBatis-Plus 3.5.5 + Druid', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, '安全与认证:', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'Spring Security 6.1 + JWT', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, 'BCrypt密码加密 + @PreAuthorize', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, 'AI服务:', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'DeepSeek API via RestTemplate\n提示词工程 + 降级策略', font_size=13, color=C_GRAY, space_before=6)

# Deployment & HarmonyOS
add_rect(s, 8.9, 1.5, 4.0, 0.45, C_DARK)
tb = add_textbox(s, 9.1, 1.53, 3.6, 0.4)
set_text(tb, '☁ 部署与鸿蒙', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 9.1, 2.1, 3.8, 4.5)
set_text(tb, '生产部署:', font_size=15, bold=True, color=C_DARK)
add_para(tb, '阿里云ECS CentOS 7.9', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, 'MySQL 8.0 + Druid连接池', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '前后端一体化Fat JAR(64MB)', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, 'nohup java -jar 后台运行', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, '鸿蒙端:', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'ArkTS + ArkUI + WebView', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, '@kit.ArkWeb Web组件', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '加载云端部署的Web页面', font_size=13, color=C_GRAY, space_before=4)

# Bottom stats bar
y_bottom = 6.5
add_rect(s, 0.3, y_bottom, 3.2, 0.7, C_BLUE_LT)
tb = add_textbox(s, 0.5, y_bottom+0.12, 2.8, 0.5)
set_text(tb, '19个 Controller', font_size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_para(tb, '后端API接口', font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, space_before=2)

add_rect(s, 3.7, y_bottom, 3.2, 0.7, C_BLUE_LT)
tb = add_textbox(s, 3.9, y_bottom+0.12, 2.8, 0.5)
set_text(tb, '27个 Vue页面', font_size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_para(tb, '前端组件', font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, space_before=2)

add_rect(s, 7.1, y_bottom, 3.0, 0.7, C_BLUE_LT)
tb = add_textbox(s, 7.3, y_bottom+0.12, 2.6, 0.5)
set_text(tb, '12张数据表', font_size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_para(tb, 'MySQL 8.0', font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, space_before=2)

add_rect(s, 10.3, y_bottom, 2.7, 0.7, C_BLUE_LT)
tb = add_textbox(s, 10.5, y_bottom+0.12, 2.3, 0.5)
set_text(tb, '3个前端', font_size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_para(tb, 'Web+PC+鸿蒙', font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, space_before=2)

# ===== Slide 5: Requirements =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '三、需求分析与功能模块', font_size=32, bold=True, color=C_WHITE)

# User side
add_rect(s, 0.3, 1.5, 6.0, 0.45, C_DARK)
tb = add_textbox(s, 0.5, 1.53, 5.6, 0.4)
set_text(tb, '👤 用户端功能 (10个模块)', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 0.5, 2.1, 5.8, 4.5)
modules = [
    ('注册/登录', 'JWT Token + BCrypt + localStorage持久化'),
    ('商品首页', 'Banner轮播 + 分类图标 + 热搜推荐'),
    ('商品搜索', '关键词模糊匹配(名/副标题/描述)'),
    ('商品详情', '移动端Vant + PC端京东双栏自适应'),
    ('购物车', 'Pinia状态管理 + 全选结算'),
    ('订单管理', '5状态流转(待付→待发→待收→完成/取消)'),
    ('收货地址', 'GPS/IP双重定位自动填充省市区'),
    ('个人中心', '卡片布局 + 订单统计 + 摄像头拍照'),
    ('AI客服', 'DeepSeek对话 + 商品推荐卡片(11品类)'),
    ('收藏管理', '收藏/取消 + 列表查看'),
]
for title, desc in modules:
    add_para(tb, f'▸ {title}：{desc}', font_size=13, color=C_DARK, space_before=8)

# Admin side
add_rect(s, 6.8, 1.5, 6.2, 0.45, C_DARK)
tb = add_textbox(s, 7.0, 1.53, 5.8, 0.4)
set_text(tb, '🔧 管理端功能 (10个模块)', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 7.0, 2.1, 5.8, 4.5)
admin_modules = [
    ('仪表盘', '四维统计 + 最近订单 + 热销Top10'),
    ('用户管理', '搜索筛选 + CRUD + Excel导出'),
    ('分类管理', '一级/二级展开式表格CRUD'),
    ('商品管理', '三级导航 + 图片上传 + 规格管理'),
    ('订单管理', '筛选 + 详情弹窗(明细表) + 一键发货'),
    ('Banner管理', 'CRUD + 图片上传 + 时间排期'),
    ('评价管理', '审核 + 回复 + 显示/隐藏 + 删除'),
    ('客服消息', '用户对话查看 + 手动回复'),
    ('系统设置', '修改密码(BCrypt) + 关于信息'),
    ('移动端适配', '全局响应式(表格/弹窗/分页)'),
]
for title, desc in admin_modules:
    add_para(tb, f'▸ {title}：{desc}', font_size=13, color=C_DARK, space_before=8)

# ===== Slide 6: Database Design =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '四、数据库设计 (12张核心数据表)', font_size=32, bold=True, color=C_WHITE)

# Core tables
add_rect(s, 0.3, 1.5, 6.0, 0.4, C_DARK)
tb = add_textbox(s, 0.5, 1.52, 5.6, 0.35)
set_text(tb, '📋 核心业务表', font_size=18, bold=True, color=C_WHITE)

tables_data = [
    ('user', '15字段', '用户名/BCrypt密码/头像/角色/状态'),
    ('category', '9字段', '二级分类/base64图标/排序/状态'),
    ('product', '19字段', '名称/价格/库存/销量/热销/新品'),
    ('product_spec', '8字段', 'SKU规格(颜色/尺码)/独立价格'),
    ('order_info', '16字段', '订单号/金额/5状态流转/支付'),
    ('order_item', '11字段', '商品快照(防历史数据变更)'),
]
y = 2.1
for name, cols, desc in tables_data:
    add_rect(s, 0.3, y, 1.3, 0.35, C_PRIMARY)
    tb = add_textbox(s, 0.4, y+0.03, 1.1, 0.3)
    set_text(tb, name, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb = add_textbox(s, 1.8, y+0.03, 4.3, 0.3)
    set_text(tb, f'{desc} ({cols})', font_size=12, color=C_DARK)
    y += 0.5

# Extended tables
add_rect(s, 6.8, 1.5, 6.2, 0.4, C_DARK)
tb = add_textbox(s, 7.0, 1.52, 5.8, 0.35)
set_text(tb, '📌 扩展业务表', font_size=18, bold=True, color=C_WHITE)

ext_tables = [
    ('address', '10字段', '收货人/省市区/GPS定位/默认'),
    ('banner', '11字段', '轮播图/图片/链接/排期'),
    ('favorite', '4字段', '用户收藏商品'),
    ('review', '13字段', '评分/内容/图片/管理员回复'),
    ('cart', '7字段', '购物车(服务端)'),
    ('customer_service', '8字段', '客服消息(sender_type)'),
]
y = 2.1
for name, cols, desc in ext_tables:
    add_rect(s, 6.8, y, 1.3, 0.35, C_GREEN)
    tb = add_textbox(s, 6.9, y+0.03, 1.1, 0.3)
    set_text(tb, name, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb = add_textbox(s, 8.3, y+0.03, 4.5, 0.3)
    set_text(tb, f'{desc} ({cols})', font_size=12, color=C_DARK)
    y += 0.5

# Design features
add_rect(s, 0.3, 5.5, 12.7, 1.5, RGBColor(0xF0,0xFD,0xF4))
tb = add_textbox(s, 0.6, 5.65, 12, 1.2)
set_text(tb, '数据库设计特点', font_size=16, bold=True, color=C_GREEN)
add_para(tb, '• 所有表均支持逻辑删除(deleted字段) — MyBatis-Plus @TableLogic自动过滤    • 订单明细快照机制 — 存储下单时商品名称/图片/价格, 防止后续商品变更影响历史订单', font_size=13, color=C_DARK, space_before=10)
add_para(tb, '• 雪花算法生成订单号(IdUtil.getSnowflakeNextIdStr) — 分布式唯一ID    • LambdaQueryWrapper参数化查询 — 防止SQL注入    • 初始数据: 40件商品 | 12条订单 | 24条评价 | 7个一级分类', font_size=13, color=C_DARK, space_before=8)

# ===== Slide 7-8: User Features =====
for slide_num, slide_title, features in [
    (7, '五、核心功能演示 - 用户端 (1/2)', [
        ('📱 商品首页', 'Banner轮播(db图片)→分类图标导航→热门/新品推荐→搜索栏\n响应式产品卡片(悬停上浮+阴影)'),
        ('🔍 商品分类与搜索', '一级分类侧边栏 + 二级Tab + 响应式商品网格(3-5列)\n图片加载失败自动显示SVG渐变占位图(12色)'),
        ('🛍 商品详情(双端)', '移动端:轮播图+价格区+规格弹窗+ActionBar\nPC端:京东风格左右双栏+大图切换+Tab详情/评价'),
        ('🛒 购物车', 'Pinia+localStorage持久化 | 全选/数量/删除 | 实时总金额\n提交订单→结算页(地址+商品清单+备注+支付)'),
    ]),
    (8, '五、核心功能演示 - 用户端 (2/2)', [
        ('📋 订单管理', '5状态Tab(全部/待付款/待发货/待收货/已完成)\n订单卡片+操作按钮(取消/支付/确认收货/详情)'),
        ('📍 收货地址', 'GPS/IP双重定位自动填充省市区(ipsapi.co→api.ip.sb回退)\nCRUD + 默认地址 + Vant AddressList组件'),
        ('👤 个人中心', '渐变卡片头像区+订单统计(API实时)+宫格服务入口\n头像上传:相册选择 + 摄像头拍照(capture="camera")'),
        ('🤖 AI客服', 'DeepSeek大模型实时对话+商品推荐(按关键词+品类)\n快捷问题标签→AI解析→推荐卡片(名称/价格/图片/链接)'),
    ]),
]:
    s = add_blank_slide()
    add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
    add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
    tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
    set_text(tb, slide_title, font_size=30, bold=True, color=C_WHITE)

    positions = [(0.3, 1.6, 6.2, 2.6), (6.8, 1.6, 6.2, 2.6), (0.3, 4.4, 6.2, 2.6), (6.8, 4.4, 6.2, 2.6)]
    for (lx, ly, lw, lh), (title, desc) in zip(positions, features):
        add_rect(s, lx, ly, lw, lh, C_LIGHT)
        add_rect(s, lx, ly, lw, 0.45, C_PRIMARY)
        tb = add_textbox(s, lx+0.2, ly+0.06, lw-0.4, 0.35)
        set_text(tb, title, font_size=15, bold=True, color=C_WHITE)
        tb = add_textbox(s, lx+0.2, ly+0.6, lw-0.4, lh-0.8)
        set_text(tb, desc.replace('\n', '\n'), font_size=12, color=C_DARK)

# ===== Slide 9-10: Admin Features =====
for slide_num, slide_title, features in [
    (9, '五、核心功能演示 - 管理端 (1/2)', [
        ('📊 仪表盘', '四个统计卡片(用户/商品/订单/销售额)\n最近10条订单表格 + 热销Top10\nElement Plus Icon彩色图标'),
        ('👥 用户管理', 'keyword/role/status筛选 + 分页\nCRUD弹窗(所有字段含头像上传)\nBCrypt加密 + Excel导出'),
        ('📂 分类管理', '一级/二级展开式表格(el-table expand)\n名称/排序/状态CRUD + 确认删除\n级联删除子分类和商品'),
        ('📦 商品管理', '三级导航: 分类→子分类→商品列表\nkeyword/status筛选 + 多图管理\n上下架/热销/新品标签 + Excel导出'),
    ]),
    (10, '五、核心功能演示 - 管理端 (2/2)', [
        ('📋 订单管理', 'orderNo/status筛选 + 分页列表\n详情Dialog:订单信息+商品明细表\n待发货→发货按钮(status流转+时间记录)'),
        ('🎨 Banner管理', 'CRUD + 图片上传(UUID命名+日期分目录)\n标题/链接/位置/排序/排期设置\nWebMvcConfigurer映射uploads URL'),
        ('⭐ 评价管理', 'productId/status筛选 + 分页\n回复Dialog + 显示/隐藏切换\nrate星级评分渲染'),
        ('⚙ 系统设置', '原密码验证 + 新密码BCrypt更新\n移动端全局响应式适配\n表格/弹窗/分页/抽屉自适应'),
    ]),
]:
    s = add_blank_slide()
    add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
    add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
    tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
    set_text(tb, slide_title, font_size=30, bold=True, color=C_WHITE)

    positions = [(0.3, 1.6, 6.2, 2.6), (6.8, 1.6, 6.2, 2.6), (0.3, 4.4, 6.2, 2.6), (6.8, 4.4, 6.2, 2.6)]
    for (lx, ly, lw, lh), (title, desc) in zip(positions, features):
        add_rect(s, lx, ly, lw, lh, C_LIGHT)
        add_rect(s, lx, ly, lw, 0.45, C_DARK)
        tb = add_textbox(s, lx+0.2, ly+0.06, lw-0.4, 0.35)
        set_text(tb, title, font_size=15, bold=True, color=C_WHITE)
        tb = add_textbox(s, lx+0.2, ly+0.6, lw-0.4, lh-0.8)
        set_text(tb, desc.replace('\n', '\n'), font_size=12, color=C_DARK)

# ===== Slide 11: AI Innovation =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11.5, 1)
set_text(tb, '六、AI智能客服创新点 ⭐ 核心亮点', font_size=30, bold=True, color=C_WHITE)
add_para(tb, 'DeepSeek大语言模型 + 检索增强生成(RAG) + 关联词映射引擎', font_size=16, color=RGBColor(0xBF,0xDB,0xFE), space_before=4)

# Flow chart
steps = [
    ('1', '用户发送消息\n"我想要饮料"', C_PRIMARY),
    ('2', '提取关键词\n+关联词映射\n(11品类60+词)', C_DARK),
    ('3', '搜索DB商品\n(上架+有库存\n销量降序Top10)', C_GREEN),
    ('4', '构建Prompt\n(商品列表+\n对话历史)', C_ORANGE),
    ('5', 'DeepSeek API\n返回推荐\n[RECOMMEND]', C_RED),
    ('6', '解析+限3个\n商品卡片嵌入\n聊天流', C_PRIMARY),
]
x = 0.5
for num, desc, color in steps:
    add_rect(s, x, 1.8, 1.8, 1.5, color)
    add_circle(s, x+0.6, 1.65, 0.6, C_WHITE)
    ct = add_textbox(s, x+0.65, 1.7, 0.5, 0.5)
    set_text(ct, num, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    ct = add_textbox(s, x+0.1, 2.4, 1.6, 0.8)
    set_text(ct, desc.replace('\n', '\n'), font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)
    if x < 11:
        arrow_tf = add_textbox(s, x+1.8, 2.3, 0.4, 0.4)
        set_text(arrow_tf, '→', font_size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        # Arrow handled above
    x += 2.15

# Bottom details
add_rect(s, 0.5, 3.6, 6.0, 3.4, C_LIGHT)
tb = add_textbox(s, 0.7, 3.75, 5.6, 3)
set_text(tb, '技术实现', font_size=18, bold=True, color=C_DARK)
add_para(tb, '• DeepSeekService: RestTemplate调用API, System Prompt动态构建(嵌入商品列表)', font_size=13, color=C_DARK, space_before=10)
add_para(tb, '• 关联词映射: HashMap<String, String[]>覆盖11品类60+关键词', font_size=13, color=C_DARK, space_before=8)
add_para(tb, '• 推荐协议: [RECOMMEND]id1,id2,id3[/RECOMMEND]标记解析,上限3个', font_size=13, color=C_DARK, space_before=8)
add_para(tb, '• 降级策略: API调用失败→catch预设关键词回复(<1s响应)', font_size=13, color=C_DARK, space_before=8)
add_para(tb, '• 前端: 消息发送→清理Markdown符号→渲染气泡+商品卡片(图片/名称/价格/跳转箭头)', font_size=13, color=C_DARK, space_before=8)
add_para(tb, '• API: model=deepseek-chat, temperature=0.7, max_tokens=800', font_size=13, color=C_DARK, space_before=8)

add_rect(s, 7.0, 3.6, 5.8, 3.4, RGBColor(0xFE,0xFB,0xEB))
tb = add_textbox(s, 7.2, 3.75, 5.4, 3)
set_text(tb, '效果展示', font_size=18, bold=True, color=C_ORANGE)
add_para(tb, '用户: "我想要饮料"', font_size=14, bold=True, color=C_DARK, space_before=10)
add_para(tb, 'AI: "为您推荐以下饮料哦~"', font_size=13, color=C_GREEN, space_before=6)
add_para(tb, '【元气森林气泡水】无糖低卡 ¥59.90', font_size=13, color=C_PRIMARY, space_before=6)
add_para(tb, '【智利进口车厘子】JJ级 ¥199.00', font_size=13, color=C_PRIMARY, space_before=4)
add_para(tb, '→ 点击商品卡片可跳转详情页购买', font_size=12, color=C_GRAY, space_before=10)
add_para(tb, '', font_size=8, color=C_GRAY, space_before=8)
add_para(tb, '用户: "推荐一款手机"', font_size=14, bold=True, color=C_DARK, space_before=4)
add_para(tb, 'AI: 推荐iPhone 15 Pro Max / 小米14 Pro / 三星S24 Ultra', font_size=13, color=C_GREEN, space_before=6)
add_para(tb, '各附价格和购买链接', font_size=12, color=C_GRAY, space_before=4)

# ===== Slide 12: Testing =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '七、作品测试', font_size=32, bold=True, color=C_WHITE)

add_rect(s, 0.5, 1.6, 5.8, 0.45, C_GREEN)
tb = add_textbox(s, 0.7, 1.63, 5.4, 0.4)
set_text(tb, '✅ 测试结果: 37个用例 100%通过', font_size=20, bold=True, color=C_WHITE)

tb = add_textbox(s, 0.7, 2.3, 5.4, 4)
set_text(tb, '用户端 24个测试用例', font_size=16, bold=True, color=C_DARK)
test_items = [
    '注册(合法/重复用户名) → 登录(正确/错误)',
    '首页加载 → 搜索 → 分类 → 商品详情',
    'PC端京东布局 → 移动端Vant布局',
    '购物车(加入/修改/删除/全选)',
    '下单(选地址→确认→提交→支付)',
    '订单(5状态Tab切换/取消/支付)',
    '收货地址(GPS定位/编辑/默认)',
    'AI客服(泛指推荐/口语回复/降级)',
    '个人中心(统计/头像拍照/资料)',
]
for item in test_items:
    add_para(tb, f'✓ {item}', font_size=12, color=C_GRAY, space_before=10)

add_rect(s, 7.0, 1.6, 5.8, 0.45, C_GREEN)
tb = add_textbox(s, 7.2, 1.63, 5.4, 0.4)
set_text(tb, '📱 兼容性测试', font_size=20, bold=True, color=C_WHITE)

tb = add_textbox(s, 7.2, 2.3, 5.4, 4)
set_text(tb, '管理端 13个测试用例', font_size=16, bold=True, color=C_DARK)
admin_test_items = [
    '仪表盘(统计/最近订单/热销)',
    '用户管理(CRUD/搜索/导出Excel)',
    '分类管理(展开/CRUD/级联删除)',
    '商品管理(三级/上架/上下架/导出)',
    '订单管理(筛选/详情弹窗/发货)',
    'Banner管理(CRUD/图片上传)',
    '评价管理(回复/显示隐藏)',
    '移动端管理(表格/弹窗/分页)',
]
for item in admin_test_items:
    add_para(tb, f'✓ {item}', font_size=12, color=C_GRAY, space_before=10)

add_rect(s, 0.5, 5.8, 12.3, 1.2, C_LIGHT)
tb = add_textbox(s, 0.7, 5.9, 11.9, 1)
set_text(tb, '测试环境', font_size=16, bold=True, color=C_DARK)
add_para(tb, '阿里云ECS 2核4G CentOS 7.9 | MySQL 8.0 | Chrome 120+ / Edge 120+ / Safari 17+ | iPhone 14 Pro / Android Chrome | 鸿蒙模拟器 | 1920x1080 / 1366x768 | 初始数据: 40商品/12订单/24评价', font_size=12, color=C_GRAY, space_before=8)

# ===== Slide 13: Problems =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '七、项目总结 - 技术问题与解决 (8个问题全部攻克)', font_size=28, bold=True, color=C_WHITE)

problems = [
    ('Vant4组件名变更', 'GoodsAction→ActionBar', '查阅CHANGELOG迁移'),
    ('高分辨率手机误判PC', '1080px手机显示PC布局', 'UA+触摸+宽度三维检测'),
    ('Vite代理rewrite问题', '/api/路径不匹配Controller', '去rewrite+统一完整路径'),
    ('Banner时间过滤NULL', 'startTime NULL→SQL不命中', '增加IS NULL OR判断'),
    ('Linux MySQL Socket', 'localhost→socket连接失败', 'JDBC URL改127.0.0.1'),
    ('SPA路由刷新500', '/admin直接访问报错', 'SpaFilter最高优先级转发'),
    ('文件上传boundary', '手动设Content-Type缺边界', '让axios自动生成请求头'),
    ('注册密码双重BCrypt', 'Controller+Service各自加密', '去Controller重复编码'),
]
y = 1.6
for i, (name, problem, solution) in enumerate(problems):
    color = [C_RED, C_ORANGE, C_GREEN, C_PRIMARY, C_DARK, C_GRAY, C_ACCENT, C_RED][i]
    add_rect(s, 0.3, y, 0.06, 0.7, color)
    add_circle(s, 0.5, y+0.15, 0.35, color)
    ct = add_textbox(s, 0.52, y+0.17, 0.3, 0.3)
    set_text(ct, str(i+1), font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb = add_textbox(s, 1.0, y+0.02, 3.2, 0.3)
    set_text(tb, name, font_size=14, bold=True, color=C_DARK)
    tb = add_textbox(s, 1.0, y+0.32, 3.2, 0.3)
    set_text(tb, f'问题: {problem}', font_size=11, color=C_GRAY)
    tb = add_textbox(s, 4.5, y+0.02, 8.5, 0.3)
    set_text(tb, f'→ 解决: {solution}', font_size=12, color=C_DARK)
    y += 0.72

# ===== Slide 14: Summary =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '八、项目总结与展望', font_size=32, bold=True, color=C_WHITE)

# Achievements
add_rect(s, 0.3, 1.5, 6.0, 0.45, C_GREEN)
tb = add_textbox(s, 0.5, 1.53, 5.6, 0.4)
set_text(tb, '🏆 项目成果', font_size=20, bold=True, color=C_WHITE)

tb = add_textbox(s, 0.5, 2.1, 5.8, 4.5)
achievements = [
    '完成12张数据表全功能电商平台',
    '后端19个Controller + 11个Service',
    '前端27个Vue页面 + 7个API模块',
    'HarmonyOS ArkTS鸿蒙端原生应用',
    'DeepSeek AI客服(11品类智能推荐)',
    '双端自适应布局(移动Vant+PC京东)',
    'GPS/IP双重定位收货地址',
    '37个测试用例100%通过',
    '部署阿里云ECS运行正常(64MB JAR)',
    'Git 25+次提交，完整开发日志',
]
for a in achievements:
    add_para(tb, f'✦ {a}', font_size=13, color=C_DARK, space_before=8)

add_rect(s, 6.8, 1.5, 6.2, 0.45, C_ORANGE)
tb = add_textbox(s, 7.0, 1.53, 5.8, 0.4)
set_text(tb, '💡 创新点 / 不足 / 改进', font_size=18, bold=True, color=C_WHITE)

tb = add_textbox(s, 7.0, 2.1, 5.8, 4.5)
set_text(tb, '核心创新', font_size=15, bold=True, color=C_DARK)
add_para(tb, '• AI客服RAG简化实现+关联词引擎', font_size=13, color=C_DARK, space_before=6)
add_para(tb, '• 三维设备检测策略(UA+触摸+宽度)', font_size=13, color=C_DARK, space_before=4)
add_para(tb, '• GPS/IP双重定位自动填充', font_size=13, color=C_DARK, space_before=4)
add_para(tb, '', font_size=6, color=C_GRAY, space_before=6)
add_para(tb, '不足与改进方向', font_size=15, bold=True, color=C_DARK, space_before=4)
add_para(tb, '• 支付:模拟→可接入支付宝/微信SDK', font_size=13, color=C_GRAY, space_before=6)
add_para(tb, '• AI:关联词→向量数据库语义推荐', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '• 图片:base64→OSS对象存储', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '• 可引入Redis缓存+消息队列', font_size=13, color=C_GRAY, space_before=4)
add_para(tb, '• 物流可接入快递100 API', font_size=13, color=C_GRAY, space_before=4)

# Bottom stats
add_rect(s, 0.3, 6.5, 12.7, 0.7, C_LIGHT)
tb = add_textbox(s, 0.5, 6.6, 12.3, 0.5)
set_text(tb, '80+文件  |  19 Controller  |  12张表  |  27 Vue页面  |  37测试用例  |  25+ Git提交  |  3终端(Web+PC+鸿蒙)', font_size=14, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# ===== Slide 15: Thanks =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 7.5, C_PRIMARY)
add_rect(s, 0, 0, 13.333, 0.08, C_ACCENT)
add_rect(s, 0, 7.42, 13.333, 0.08, C_ACCENT)

tb = add_textbox(s, 1.5, 1.5, 10.3, 2)
set_text(tb, '感谢聆听', font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_para(tb, '请各位老师批评指正', font_size=28, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER, space_before=20)

add_rect(s, 4.5, 3.8, 4.3, 0.01, C_WHITE)

tb = add_textbox(s, 3, 4.2, 7.3, 2.5)
set_text(tb, '"柒玖商店" - 基于SpringBoot多前端多平台的商城系统', font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
add_para(tb, '', font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=16)
add_para(tb, '姓   名：李  鑫', font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=8)
add_para(tb, '学   号：202330302152', font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=4)
add_para(tb, '专   业：软件技术  |  班 级：2307', font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=4)
add_para(tb, '指导教师：彭德宇  |  2026年5月', font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=4)

# ===== Slide 16: Appendix - Code =====
s = add_blank_slide()
add_rect(s, 0, 0, 13.333, 0.06, C_ACCENT)
add_rect(s, 0, 0, 13.333, 1.2, C_PRIMARY)
tb = add_textbox(s, 0.8, 0.25, 11, 0.7)
set_text(tb, '附录：关键技术代码示例', font_size=28, bold=True, color=C_WHITE)

# JWT code
add_rect(s, 0.3, 1.5, 6.2, 2.8, RGBColor(0x1E,0x29,0x3B))
tb = add_textbox(s, 0.5, 1.6, 5.8, 2.5)
set_text(tb, 'JWT认证过滤器 (JwtAuthenticationFilter)', font_size=14, bold=True, color=C_GREEN)
code_lines = [
    'String token = authHeader.substring(7);',
    'Claims claims = jwtTokenUtil.parseToken(token);',
    'Long userId = claims.get("userId", Long.class);',
    'Integer role = claims.get("role", Integer.class);',
    'String roleName = role >= 1 ? "ROLE_ADMIN" : "ROLE_USER";',
    'UsernamePasswordAuthenticationToken authToken =',
    '  new UsernamePasswordAuthenticationToken(',
    '    userId, null,',
    '    List.of(new SimpleGrantedAuthority(roleName))',
    '  );',
    'SecurityContextHolder.getContext()',
    '  .setAuthentication(authToken);',
]
for cl in code_lines:
    add_para(tb, cl, font_size=10, color=RGBColor(0xE2,0xE8,0xF0), space_before=4)

# Device detection code
add_rect(s, 6.8, 1.5, 6.2, 2.8, RGBColor(0x1E,0x29,0x3B))
tb = add_textbox(s, 7.0, 1.6, 5.8, 2.5)
set_text(tb, '三维设备检测 (responsive.js)', font_size=14, bold=True, color=C_GREEN)
code_lines2 = [
    'const ua = navigator.userAgent || ""',
    'const isMobileUA = /Android|iPhone|iPad',
    '  |iPod|Mobile/i.test(ua)',
    'const hasTouch = "ontouchstart" in window',
    '  || navigator.maxTouchPoints > 0',
    '',
    'function detectMobile() {',
    '  if (isMobileUA) return true',
    '  if (hasTouch && innerWidth < 1024)',
    '    return true',
    '  if (!hasTouch && innerWidth < 600)',
    '    return true',
    '  return false',
    '}',
]
for cl in code_lines2:
    add_para(tb, cl, font_size=10, color=RGBColor(0xE2,0xE8,0xF0), space_before=4)

# HarmonyOS code
add_rect(s, 0.3, 4.6, 6.2, 2.5, RGBColor(0x1E,0x29,0x3B))
tb = add_textbox(s, 0.5, 4.7, 5.8, 2.2)
set_text(tb, 'HarmonyOS ArkTS入口 (Index.ets)', font_size=14, bold=True, color=C_GREEN)
code_lines3 = [
    'import { webview } from "@kit.ArkWeb"',
    '',
    '@Entry @Component',
    'struct Index {',
    '  build() {',
    '    Column() {',
    '      Web({ src: "http://47.100.214.45:',
    '        8080/", controller:',
    '        new webview.WebviewController()',
    '      }).width("100%").height("100%")',
    '      .domStorageAccess(true)',
    '      .javaScriptAccess(true)',
    '      .mixedMode(MixedMode.All)',
    '    }',
    '  }',
    '}',
]
for cl in code_lines3:
    add_para(tb, cl, font_size=10, color=RGBColor(0xE2,0xE8,0xF0), space_before=4)

# AI code
add_rect(s, 6.8, 4.6, 6.2, 2.5, RGBColor(0x1E,0x29,0x3B))
tb = add_textbox(s, 7.0, 4.7, 5.8, 2.2)
set_text(tb, 'DeepSeek AI关联词映射 (摘录)', font_size=14, bold=True, color=C_GREEN)
code_lines4 = [
    'RELATED_WORDS.put("饮料",',
    '  new String[]{"气泡水","可乐","橙汁",',
    '  "果汁","水","茶","咖啡","奶","饮"});',
    'RELATED_WORDS.put("手机",',
    '  new String[]{"iPhone","华为","小米",',
    '  "OPPO","vivo","三星","手机"});',
    'RELATED_WORDS.put("衣服",',
    '  new String[]{"羽绒","卫衣","大衣",',
    '  "连衣裙","西装","冲锋衣","衣服",',
    '  "CK","优衣库","Nike","ZARA"});',
    '// 共11个品类 60+关键词映射',
    '// 覆盖: 饮料/手机/电脑/耳机/零食/',
    '//   衣服/鞋/护肤/家电/酒/手表',
]
for cl in code_lines4:
    add_para(tb, cl, font_size=10, color=RGBColor(0xE2,0xE8,0xF0), space_before=4)

# Save
output_path = r"D:\QJShop\柒玖商店_毕业设计答辩_v2.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
