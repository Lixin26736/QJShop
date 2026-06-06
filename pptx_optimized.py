"""毕业设计答辩PPT - 按六大维度全面优化版"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs=Presentation()
prs.slide_width=Inches(13.333)
prs.slide_height=Inches(7.5)

# === 配色(深蓝主色+浅灰辅助+橙色强调) ===
P=RGBColor(0x1E,0x40,0x8E)   # 深蓝主色
PL=RGBColor(0x25,0x63,0xEB)  # 蓝色次
D=RGBColor(0x1E,0x29,0x3B)   # 深色文字
G=RGBColor(0x64,0x74,0x8B)   # 灰色辅助
L=RGBColor(0xF8,0xFA,0xFC)   # 浅灰背景
O=RGBColor(0xEA,0x58,0x0C)   # 橙色强调
E=RGBColor(0x05,0x96,0x69)   # 绿色
W=RGBColor(0xFF,0xFF,0xFF)
BL=RGBColor(0xE2,0xE8,0xF0)  # 边框色
GD=RGBColor(0x24,0x24,0x3E)  # 代码背景
GT=RGBColor(0xE2,0xE8,0xF0)  # 代码文字
INFO=("李鑫","202330302152","彭德宇")  # 页脚信息

M=0.6  # 1.5cm安全边距

def sl():return prs.slides.add_slide(prs.slide_layouts[6])
def tx(s,l,t,w,h):return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
def S(tf,x,fs=14,b=False,c=D,a=PP_ALIGN.LEFT):
    tf.word_wrap=1;p=tf.paragraphs[0];p.text=x;p.font.name='Microsoft YaHei'
    p.font.size=Pt(fs);p.font.bold=b;p.font.color.rgb=c;p.alignment=a;return p
def A(tf,x,fs=12,c=G,a=PP_ALIGN.LEFT,sb=4):
    p=tf.add_paragraph();p.text=x;p.font.name='Microsoft YaHei';p.font.size=Pt(fs)
    p.font.color.rgb=c;p.alignment=a;p.space_before=Pt(sb);return p
def bx(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background();return sh
def ft(s):
    """页脚: 姓名 学号 页码"""
    n=len(prs.slides)
    t=tx(s,M,7.05,11.7,0.3)
    S(t,f'{INFO[0]} | {INFO[1]} | 指导教师: {INFO[2]} | {n}/31',fs=9,c=G,a=PP_ALIGN.RIGHT)
    # 顶部细线
    bx(s,M,7.02,11.7,0.01,BL)

def hdr(s,ch,title,sub=None):
    """标准页头: 蓝色左竖条+章节+标题"""
    bx(s,M,M,0.06,0.5,P)
    t=tx(s,M+0.2,M-0.05,11,0.6)
    S(t,ch,fs=10,c=P,b=True)
    p=t.add_paragraph();p.text=title;p.font.name='Microsoft YaHei'
    p.font.size=Pt(26);p.font.bold=True;p.font.color.rgb=D;p.alignment=PP_ALIGN.LEFT
    if sub:
        p2=t.add_paragraph();p2.text=sub;p2.font.name='Microsoft YaHei'
        p2.font.size=Pt(13);p2.font.color.rgb=G;p2.alignment=PP_ALIGN.LEFT;p2.space_before=Pt(2)
    bx(s,M,M+0.75,11.7,0.01,BL)
    ft(s)

def card(s,lx,ly,w,h,title,lines,color=P):
    """圆角卡片"""
    bx(s,lx,ly,w,h,W)
    bx(s,lx,ly,w,0.04,color)  # 顶部颜色条
    t=tx(s,lx+0.15,ly+0.12,w-0.3,h-0.25)
    S(t,title,fs=13,b=True,c=color)
    for ln in lines:A(t,ln,fs=11,c=G,sb=3)

def stat_card(s,lx,ly,w,h,num,label,color=P):
    """数字统计卡片"""
    bx(s,lx,ly,w,h,color)
    t=tx(s,lx+0.1,ly+h*0.15,w-0.2,h*0.5)
    S(t,str(num),fs=28,b=True,c=W,a=PP_ALIGN.CENTER)
    A(t,label,fs=10,c=W,a=PP_ALIGN.CENTER,sb=4)

def code(s,lx,ly,w,h,title,lines,color=E):
    """代码块(精简核心代码)"""
    bx(s,lx,ly,w,h,GD)
    t=tx(s,lx+0.12,ly+0.06,w-0.24,h-0.12)
    S(t,title,fs=10,b=True,c=color)
    for ln in lines:
        p=t.add_paragraph();p.text=ln;p.font.name='Consolas';p.font.size=Pt(9)
        p.font.color.rgb=GT;p.alignment=PP_ALIGN.LEFT;p.space_before=Pt(2)

def img_placeholder(s,lx,ly,w,h,label):
    """图片占位"""
    bx(s,lx,ly,w,h,BL)
    t=tx(s,lx+0.1,ly+h*0.35,w-0.2,0.4)
    S(t,label,fs=10,c=G,a=PP_ALIGN.CENTER)

# ====== P01 封面 ======
s=sl()
bx(s,0,0,13.333,7.5,P)  # 深蓝全屏
bx(s,3,M+1,7.3,0.01,W)   # 装饰线
t=tx(s,2,M+1.3,9.3,1.8)
S(t,'柒玖商店',fs=48,b=True,c=W,a=PP_ALIGN.CENTER)
A(t,'基于 SpringBoot 多前端多平台的商城系统',fs=22,c=W,a=PP_ALIGN.CENTER,sb=8)
t=tx(s,3.5,M+3.5,6.3,1.2)
S(t,'毕业设计答辩',fs=20,b=True,c=W,a=PP_ALIGN.CENTER)
A(t,'李 鑫 | 202330302152 | 软件技术 2307 班 | 彭德宇 老师 | 2026 年 5 月',fs=13,c=W,a=PP_ALIGN.CENTER,sb=14)

# ====== P02 答辩大纲 ======
s=sl()
hdr(s,'CONTENTS','答辩大纲')
toc=[('01','项目背景\n与选题意义'),('02','技术架构\n概述'),('03','需求分析与\n功能模块'),('04','数据库\n设计'),('05','用户端功能\n演示(上)'),('06','用户端功能\n演示(下)'),('07','管理端功能\n演示(上)'),('08','管理端功能\n演示(下)'),('09','AI 智能客服\n(核心亮点)'),('10','作品测试与\n项目总结')]
for i,(num,title) in enumerate(toc):
    row=i//5;col=i%5
    lx=M+col*2.35;ly=1.4+row*2.7
    bx(s,lx,ly,2.15,2.4,W)
    bx(s,lx,ly,2.15,0.06,P)
    t=tx(s,lx+0.1,ly+0.15,1.95,0.5)
    S(t,num,fs=22,b=True,c=P,a=PP_ALIGN.CENTER)
    t=tx(s,lx+0.1,ly+1.2,1.95,1)
    S(t,title.replace('\n','\n'),fs=14,c=D,a=PP_ALIGN.CENTER)

# ====== P03 项目背景 ======
s=sl()
hdr(s,'01 项目背景','电商发展现状','2025年中国网络零售市场超过15万亿元，电商与AI深度融合')
# 左: 数据卡片列
for i,(v,lb) in enumerate([('15万亿+','2025年网络零售规模'),('30%+','占社消零售总额比重'),('7x24h','全天候突破时空限制')]):
    stat_card(s,M,1.2+i*1.7,3.5,1.5,v,lb,[P,D,E][i])
# 右: 趋势要点
t=tx(s,4.5,1.2,8.3,2)
S(t,'趋势驱动因素',fs=18,b=True,c=D)
for ln in ['线上消费习惯全面养成，移动支付渗透率超90%','AI大模型(DeepSeek等)为电商智能化带来全新可能','鸿蒙生态扩展，多终端适配成为应用标配','中小商家数字化转型需求旺盛，但技术门槛高']:
    A(t,ln,fs=14,c=G,sb=10)
# 底部总结
bx(s,M,6.4,11.7,0.6,L)
t=tx(s,M+0.3,6.45,11.1,0.5)
S(t,'核心洞察: 电商持续增长 | AI+鸿蒙新机遇 | 中小商家需要低成本智能化方案',fs=13,b=True,c=O,a=PP_ALIGN.CENTER)

# ====== P04 选题意义 ======
s=sl()
hdr(s,'01 项目背景','选题意义','四大维度阐述课题价值')
for i,(title,lines,c) in enumerate([
    ('技术整合实践',['SpringBoot+Vue3+鸿蒙全栈集成','JWT认证+MyBatis-Plus高效持久层','Vant移动+ElementPlus管理双UI'],P),
    ('AI 创新落地',['DeepSeek大模型电商场景应用','RAG简化实现+关联词推荐引擎','API失败自动降级(<1s预设回复)'],O),
    ('多端适配工程',['移动Web(Vant4)+PC(ElementPlus)','鸿蒙ArkTS原生WebView','三维设备检测(UA+触摸+宽度)'],E),
    ('商业实用价值',['中小商家低成本解决方案','12张表完整电商数据模型','40商品+12订单真实数据验证'],D)]):
    lx=M+col*3.05;ly=1.2+row*3
    bx(s,lx,ly,2.85,2.7,W)
    bx(s,lx,ly,2.85,0.05,c)
    t=tx(s,lx+0.15,ly+0.15,2.55,2.3)
    S(t,title,fs=16,b=True,c=c)
    for ln in lines:A(t,ln,fs=12,c=G,sb=8)

# ====== P05 技术架构全景 ======
s=sl()
hdr(s,'02 技术架构','系统全景','三层分层架构: 前端展示层 + 后端服务层 + 基础设施层')
for i,(layer,tech,desc,color) in enumerate([
    ('前端展示层','Vue3 + Vite + Vant4 + ElementPlus + ArkTS','27个Vue页面 | 7个API模块 | PC/移动/鸿蒙三端覆盖',P),
    ('后端服务层','SpringBoot3.2 + MyBatis-Plus + Security + JWT','19个Controller | 11个Service | DeepSeek AI服务',D),
    ('基础设施层','MySQL8.0 + Druid + CentOS7.9','12张数据表 | FatJAR 64MB | java -jar 一键部署',E)]):
    bx(s,M,1.2+i*1.8,11.7,1.6,L)
    bx(s,M,1.2+i*1.8,0.08,1.6,color)
    t=tx(s,M+0.3,1.3+i*1.8,5,1.3)
    S(t,layer,fs=18,b=True,c=color)
    A(t,tech,fs=14,c=D,sb=4)
    A(t,desc,fs=12,c=G,sb=4)
bx(s,M,6.6,11.7,0.4,L)
t=tx(s,M+0.3,6.62,11.1,0.35)
S(t,'架构特点: 前后端分离 | RESTful API JSON通信 | JWT无状态认证 | 前后端一体化FatJAR部署',fs=12,b=True,c=P,a=PP_ALIGN.CENTER)

# ====== P06 前端技术栈 ======
s=sl()
hdr(s,'02 技术架构','前端技术栈详解','三端覆盖: 移动端Vant4 / PC端ElementPlus / 鸿蒙端ArkTS')
for i,(title,lines,c) in enumerate([
    ('移动端用户商城 (Vant4)',['组件: NavBar/Tabbar/Swipe/Sidebar/ActionBar','状态: Pinia cartStore+userStore','适配: 320px-768px触摸优化'],P),
    ('PC管理后台 (ElementPlus)',['布局: el-container/aside/header/main','表格: el-table+el-pagination分页','移动端: CSS媒体查询全局响应式'],D),
    ('鸿蒙原生端 (ArkTS)',['语言: ArkTS 1.2.1 + ArkUI框架','组件: @kit.ArkWeb WebView加载','配置: domStorageAccess+javaScriptAccess'],E)]):
    card(s,M+i*4.1,1.2,3.9,2.2,title,lines,c)
# 设备检测代码
code(s,M,3.7,5.8,2.5,'设备检测策略 (responsive.js)',['const ua = navigator.userAgent','const isMobileUA = /Android|iPhone|iPad|Mobile/i.test(ua)','const hasTouch = "ontouchstart" in window','function detectMobile() {','  if (isMobileUA) return true','  if (hasTouch && w<1024) return true','  if (!hasTouch && w<600) return true','}'])
code(s,6.8,3.7,5.8,2.5,'首页并行加载 (Home.vue)',['const [categories,hot,new_p,banners]=','  await Promise.all([','    categoryApi.getCategoryList(),','    productApi.getHotProducts(),','    productApi.getNewProducts(),','    bannerApi.list("home").catch(()=>[])','  ])  // 四接口并行,首屏加载优化'])

# ====== P07 后端技术栈 ======
s=sl()
hdr(s,'02 技术架构','后端技术栈详解','SpringBoot3.2.5 + MyBatis-Plus3.5.5 + Security6.1 + JWT')
for i,(title,lines,c) in enumerate([
    ('认证授权模块',['JwtAuthenticationFilter 解析Token','提取userId/username/role→Authentication','@PreAuthorize("hasRole(''ADMIN'')")权限控制','SecurityFilterChain: 默认允许+注解授权'],P),
    ('持久层模块 (MyBatis-Plus)',['LambdaQueryWrapper类型安全查询','BaseMapper零XML CRUD','@TableLogic全表逻辑删除','@Select注解自定义统计SQL'],D),
    ('业务与AI模块',['11个Service+Impl业务实现','DeepSeekService AI对话+推荐','FileUploadController图片上传','ApachePOI Excel导出服务'],E)]):
    card(s,M+i*4.1,1.2,3.9,2.5,title,lines,c)
# 项目数据
for i,(v,lb,c) in enumerate([(19,'Controller',P),(11,'Service',D),(12,'Mapper',E),(27,'Vue页面',O),(37,'测试用例',P),(7,'API模块',D),(12,'数据表',E),(3,'终端',O)]):
    stat_card(s,M+(i%8)*1.5,4.2+(i//8)*1.6,1.35,1.3,v,lb,c)
bx(s,M,6.6,11.7,0.4,L)
t=tx(s,M+0.3,6.62,11.1,0.35)
S(t,'80+源文件 | 19 Controller | 11 Service | 12 Mapper | 12 Table | 27 Page | 37 Test | 3 Terminal',fs=12,b=True,c=P,a=PP_ALIGN.CENTER)

# ====== P08 用户端10模块 ======
s=sl()
hdr(s,'03 需求分析','用户端功能模块','10个核心模块覆盖完整购物流程')
mods=[('商品浏览','Banner+分类+搜索\n热门+新品推荐'),('商品详情','PC京东双栏布局\n移动Vant轮播'),('购物车','Pinia状态管理\n全选/结算'),('订单管理','5状态Tab切换\n取消/支付/收货'),('收货地址','GPS/IP双重定位\nCRUD+默认标记'),
      ('个人中心','卡片布局+统计\n摄像头拍照头像'),('收藏管理','收藏/取消+列表\n跳转商品详情'),('AI客服','DeepSeek对话\n商品推荐卡片'),('注册登录','JWT+BCrypt\n表单校验+守卫'),('下单支付','地址→确认→提交\n模拟支付弹窗')]
for i,(title,desc) in enumerate(mods):
    row=i//5;col=i%5
    lx=M+col*2.35;ly=1.2+row*3.05
    bx(s,lx,ly,2.15,2.8,W)
    bx(s,lx,ly,2.15,0.05,P)
    t=tx(s,lx+0.1,ly+0.15,1.95,2.5);S(t,title,fs=14,b=True,c=P,a=PP_ALIGN.CENTER)
    for ln in desc.split('\n'):A(t,ln,fs=11,c=G,a=PP_ALIGN.CENTER,sb=6)

# ====== P09 管理端10模块 ======
s=sl()
hdr(s,'03 需求分析','管理端功能模块','10个管理模块覆盖平台全维度数据管理')
for i,(title,desc) in enumerate([
    ('仪表盘','四维统计卡片\n订单+热销表格'),('用户管理','搜索筛选CRUD\nExcel导出'),('分类管理','一级/二级展开\n级联删除'),('商品管理','三级导航筛选\n上下架/导出'),('订单管理','筛选+详情弹窗\n一键发货'),
    ('Banner管理','CRUD+上传\n时间排期设置'),('评价管理','回复+显示隐藏\n评分展示'),('客服消息','对话查看\n手动回复'),('系统设置','密码修改\n关于信息'),('移动适配','表格/弹窗/分页\n抽屉响应式')]):
    row=i//5;col=i%5
    lx=M+col*2.35;ly=1.2+row*3.05
    bx(s,lx,ly,2.15,2.8,W)
    bx(s,lx,ly,2.15,0.05,D)
    t=tx(s,lx+0.1,ly+0.15,1.95,2.5);S(t,title,fs=14,b=True,c=D,a=PP_ALIGN.CENTER)
    for ln in desc.split('\n'):A(t,ln,fs=11,c=G,a=PP_ALIGN.CENTER,sb=6)

# ====== P10 数据库总览 ======
s=sl()
hdr(s,'04 数据库设计','12张核心数据表总览','覆盖商品/订单/用户/内容四大领域')
for i,(name,desc,c) in enumerate([
    ('user','15字段·用户/角色/状态',P),('product','19字段·价格/库存/热销',P),('category','9字段·二级/base64图标',P),('order_info','16字段·5状态流转',P),
    ('order_item','11字段·商品快照',D),('product_spec','8字段·SKU规格',D),('address','10字段·GPS定位',D),('banner','11字段·轮播排期',D),
    ('review','13字段·评分回复',E),('favorite','4字段·用户收藏',E),('cart','7字段·购物车',E),('customer_service','8字段·客服消息',E)]):
    row=i//4;col=i%4
    lx=M+col*3;ly=1.2+row*2
    bx(s,lx,ly,2.8,1.7,L)
    bx(s,lx,ly,2.8,0.05,c)
    t=tx(s,lx+0.1,ly+0.15,2.6,0.5);S(t,name,fs=15,b=True,c=c)
    A(t,desc,fs=11,c=G,sb=4)

# ====== P11 数据库特点 ======
s=sl()
hdr(s,'04 数据库设计','设计特点与初始数据','逻辑删除 / 快照机制 / 雪花ID / 参数化查询')
card(s,M,1.2,5.8,3,'设计特点',['逻辑删除: 12张表全部deleted字段, @TableLogic自动过滤','订单快照: order_item存储下单时商品名/价格/图片','雪花ID: Hutool IdUtil.getSnowflakeNextIdStr()','BCrypt密码: Spring Security BCryptPasswordEncoder','参数化查询: LambdaQueryWrapper防SQL注入'],P)
card(s,7,1.2,5.8,3,'初始数据统计',['商品: 40件(手机/电脑/家电/服装/美妆/食品/运动等)','订单: 12条(覆盖5种状态流转)','评价: 24条(含评分/图片/管理员回复)','规格: 57条 | 地址: 15条 | Banner: 5条(Unsplash)','用户: 8名(含管理员) | 分类: 7一级+30二级'],D)
code(s,M,4.5,5.8,2.2,'MyBatis-Plus逻辑删除配置',['@TableLogic','private Integer deleted;','// deleteById → UPDATE ... SET deleted=1','// selectList → WHERE deleted=0 自动过滤'])
code(s,7,4.5,5.8,2.2,'雪花算法生成订单号',['import cn.hutool.core.util.IdUtil;','String orderNo = IdUtil.getSnowflakeNextIdStr();','// 分布式唯一ID, 无需数据库自增','// 示例: 1847654321098765432'])

# ====== P12 首页与商品浏览 ======
s=sl()
hdr(s,'05 用户端功能(上)','首页与商品浏览','/client/home | /client/category | /client/product/:id')
card(s,M,1.2,5.8,2.5,'首页功能 (Home.vue)',['Banner轮播: VantSwipe 3s自动播放, 5条数据','分类导航: DB中base64图标, 响应式4-5列','热门商品: /api/products/hot 销量Top10','新品推荐: /api/products/new 时间Top10','搜索栏: 圆形搜索框, keyword模糊匹配'],P)
card(s,7,1.2,5.8,2.5,'分类与搜索 (Category.vue)',['一级分类: VantSidebar侧边栏点击切换','二级分类: VantTabs标签动态加载','响应式网格: CSSGrid 3列(小屏)/4列/5列(大屏)','搜索模式: 隐藏侧边栏→调用/api/products/search','占位图: SVG渐变+商品首字 12色轮换'],D)
code(s,M,4,5.8,2.4,'商品搜索后端实现',['LambdaQueryWrapper<Product> w = new LambdaQueryWrapper<>()','w.eq(Product::getStatus, 1).gt(Product::getStock, 0)','w.and(wr -> {','  for (String word : words)','    wr.or().like(Product::getName, word)','      .or().like(Product::getSubtitle, word)','      .or().like(Product::getDescription, word)','})','w.orderByDesc(Product::getSales).last("limit 10")'])
img_placeholder(s,7,4,5.8,2.4,'[此处插入: 首页截图 + 分类页截图]')

# ====== P13 商品详情双端布局 ======
s=sl()
hdr(s,'05 用户端功能(上)','商品详情 — 双端自适应布局','移动端Vant4轮播+ActionBar | PC端京东风格左右双栏')
card(s,M,1.2,5.8,2.8,'移动端布局 (Vant4)',['顶部: VantSwipe图片轮播(3s自动)','价格区: 红色¥价格+原价划线+销量/库存','规格选择: VantCell→底部Popup弹窗标签选择','底部ActionBar: 客服/收藏/购物车/加购/购买','占位图: SVG渐变+商品首字(12色)'],P)
card(s,7,1.2,5.8,2.8,'PC端布局 (京东风格)',['左侧450px: 大图+60px缩略图条点击切换','右侧价格区: ¥28px红字+原价14px划线+折扣%','规格选择: 标签直接点击切换(价格联动)','数量: VantStepper min=1 max=库存','大按钮: 加购160x48/购买160x48/收藏/客服'],D)
code(s,M,4.3,5.8,2.4,'PC端缩略图切换逻辑',['const activeImg = ref(product.mainImage)','const images = computed(() => {','  const imgs = [product.mainImage]','  // JSON.parse(detailImages) 兼容数组/逗号分隔','  if (product.detailImages) {','    const arr = JSON.parse(product.detailImages)','    arr.forEach(img => imgs.push(img))','  }','  return imgs.length ? imgs : [placeholder]','})'])
img_placeholder(s,7,4.3,5.8,2.4,'[此处插入: 移动端商品详情截图 | PC端商品详情截图 — 左右分栏对比]')

# ====== P14 购物车+订单+地址 ======
s=sl()
hdr(s,'05 用户端功能(上)','购物车 / 订单管理 / 收货地址','核心交易流程: 加购 → 下单 → 支付 → 发货 → 收货')
card(s,M,1.2,3.9,2.5,'购物车 (Cart.vue)',['Pinia cartStore集中管理','localStorage持久化','全选/数量/删除/实时总金额','提交→结算页(地址+商品+备注)'],P)
card(s,4.7,1.2,3.9,2.5,'订单管理 (OrderList)',['5状态Tab(全部/待付/待发/待收/完成)','切换Tab→resetAndLoad()','取消/支付/确认收货按钮','详情: /client/order/:id'],D)
card(s,9.4,1.2,3.9,2.5,'收货地址 (AddressList)',['VantAddressList+默认标记','GPS定位→Nominatim逆地理编码','权限拒绝→IP定位(ipapi+ipsb)','CRUD+设置默认'],E)
code(s,M,4,6.5,1.8,'订单提交核心代码',['const items = checkoutItems.map(i => ({','  productId: i.id, quantity: i.quantity','}))','const res = await orderApi.createOrder({','  addressId: selectedAddress.id, items, remark','})  // 返回: { orderId, orderNo, totalAmount }'])
code(s,7,4,5.8,1.8,'GPS/IP双重定位策略',['// 1. GPS: navigator.geolocation','// 2. 逆地理: Nominatim API','// 3. 权限拒绝→IP定位: ipapi.co/json','// 4. 备用: api.ip.sb/geoip','// 全失败→提示手动输入'])

# ====== P15 个人中心+AI客服 ======
s=sl()
hdr(s,'05 用户端功能(下)','个人中心 / 账号设置 / AI客服','/client/profile | /client/settings | /client/cs')
card(s,M,1.2,3.9,2.5,'个人中心 (Profile)',['渐变卡片头像区+昵称/手机','订单统计: 5状态数量(API实时)','服务宫格: 订单/地址/收藏/客服','退出登录(确认弹窗)'],P)
card(s,4.7,1.2,3.9,2.5,'账号设置 (Settings)',['头像: 相册选择+摄像头拍照','capture="camera" 调用设备','大小限制: 10MB→FileReader→base64','资料: 昵称/手机/邮箱/性别/生日'],D)
card(s,9.4,1.2,3.9,2.5,'AI客服 (CustomerService)',['欢迎语: 工号26793小琉','快捷问题: 5个标签一键发送','AI思考: 三点闪烁动画','回复: Markdown过滤+商品卡片'],E)
code(s,M,4,6.5,1.8,'AI客服前端核心逻辑',['const sendMessage = async () => {','  const res = await csApi.sendMessage(text)','  let reply = res.reply','    .replace(/[*#_`~>-]/g, "")  // 去符号','    .replace(/\\[RECOMMEND\\].*?\\/g, "")','  messages.push({content: reply,','    products: res.products})','}'])
img_placeholder(s,7,4,5.8,1.8,'[此处插入: 个人中心截图 | AI客服对话截图(含推荐卡片)]')

# ====== P16 管理端仪表盘+用户管理 ======
s=sl()
hdr(s,'06 管理端功能(上)','仪表盘与用户管理','/admin/dashboard | /admin/user')
card(s,M,1.2,5.8,2.8,'仪表盘 (Dashboard.vue)',['统计卡片: 4个ElementPlusCard+彩色图标','  (用户总数/商品总数/订单总数/销售总额)','最近订单: el-table最近10条(订单号/金额/状态)','热销商品: 销量Top10(商品名/销量/价格)','数据源: DashboardStatsVO统一封装返回'],P)
card(s,7,1.2,5.8,2.8,'用户管理 (User.vue)',['搜索: keyword(用户名/昵称/手机)+role+status','列表: ID/用户名/昵称/手机/邮箱/角色/状态','CRUD: Dialog表单(含头像上传→base64)','导出: ApachePOI SXSSFWorkbook→xlsx下载','密码: UserServiceImpl.update() BCrypt加密'],D)
code(s,M,4.3,5.8,1.8,'仪表盘数据统计',['// DashboardServiceImpl.stats()','Long userCount = userMapper.selectCount(wrapper)','Long productCount = productMapper.selectCount(wrapper)','BigDecimal totalSales = orderMapper.sumPayAmount()','Long todayOrders = orderMapper.countTodayOrders(start,end)'])
code(s,7,4.3,5.8,1.8,'用户Excel导出',['@GetMapping("/export")','public void export(HttpServletResponse res){','  List<User> users = userService.listAll(...)','  String[] headers={"ID","用户名","昵称",...}','  ExcelExportUtil.exportExcel(res,headers,data)','}'])

# ====== P17 管理端分类+商品管理 ======
s=sl()
hdr(s,'06 管理端功能(上)','分类管理与商品管理','/admin/category | /admin/product (三级导航)')
card(s,M,1.2,5.8,2.8,'分类管理 (Category.vue)',['一级分类: el-table + type="expand"','  展开显示二级子分类(独立子表格)','  ID/名称/排序/状态标签/编辑/删除','Dialog: 名称(input)+排序(number)+状态(select)','删除: ElMessageBox.confirm(级联提示)'],P)
card(s,7,1.2,5.8,2.8,'商品管理 (Product.vue)',['三级导航: 一级分类→查看子分类→查看商品','  el-breadcrumb面包屑路径','列表: keyword+status筛选+el-pagination','Dialog: 18字段完整表单(名称/价格/库存/等)','  主图/多详情图上传管理','导出: /api/admin/products/export'],D)
code(s,M,4.3,5.8,1.8,'商品分页查询',['LambdaQueryWrapper<Product> w = new LambdaQueryWrapper<>()','w.eq(categoryId!=null, Product::getCategoryId, id)','  .like(keyword!=null, Product::getName, keyword)','  .eq(status!=null, Product::getStatus, status)','  .orderByDesc(Product::getSortOrder)','Page<Product> p = mapper.selectPage(new Page<>(), w)'])
code(s,7,4.3,5.8,1.8,'分类子分类加载',['const getChildren = (parentId) => {','  if (!childrenCache[parentId]) {','    request.get(`/${parentId}`).then(res=>{','      childrenCache = {...childrenCache,','        [parentId]: res}  // 对象替换触发响应式','    })','  }','}'])

# ====== P18 管理端订单+Banner ======
s=sl()
hdr(s,'06 管理端功能(下)','订单管理与Banner管理','/admin/order | /admin/banner')
card(s,M,1.2,5.8,2.8,'订单管理 (Order.vue)',['搜索: orderNo+status筛选+分页','列表: ID/订单号/金额/支付方式/状态/时间','查看: Dialog弹窗→订单信息+商品明细table','发货: status=1→点击发货→status=2(待收货)','  记录delivery_time','导出: /api/admin/orders/export'],P)
card(s,7,1.2,5.8,2.8,'Banner管理 (Banner.vue)',['列表: 缩略图/标题/位置/排序/状态/时间','Dialog: 标题+图片上传+链接+排序+时间','上传: FormData→/api/upload/image','  UUID命名→uploads/日期/分目录→返回URL','  前端: uploadApi.uploadImage(file)','状态: 启用/禁用(时间排期+IS NULL处理)'],D)
code(s,M,4.3,5.8,1.8,'订单发货状态流转',['@PutMapping("/{id}/ship")','public ApiResponse<Void> ship(@PathVariable Long id){','  OrderInfo o = service.getById(id)','  if (o != null && o.getStatus() == 1) {','    o.setStatus(2);  // 待收货','    o.setDeliveryTime(LocalDateTime.now())','    service.update(o)','  }','}'])
code(s,7,4.3,5.8,1.8,'图片上传处理',['@PostMapping("/image")','public ApiResponse uploadImage(@RequestParam file){','  if (!contentType.startsWith("image/"))','    return error("只支持图片")','  String dir = now().format("yyyy/MM/dd")','  String name = UUID.randomUUID()+ext','  file.transferTo(new File(uploadDir,dir,name))','  return success("/uploads/"+dir+"/"+name)','}'])

# ====== P19 评价管理+系统设置 ======
s=sl()
hdr(s,'06 管理端功能(下)','评价管理与系统设置','/admin/review | /admin/settings | 移动端响应式适配')
card(s,M,1.2,3.9,2.5,'评价管理 (Review)',['productId+status筛选','Dialog回复(展示原评价+textarea)','保存: reply+reply_time更新','显示/隐藏切换+删除'],P)
card(s,4.7,1.2,3.9,2.5,'系统设置 (Settings)',['修改密码: 原密码+新密码+确认','BCrypt加密→清除Token重登','关于: 系统版本+技术栈版本','客服消息: 对话查看+回复'],D)
card(s,9.4,1.2,3.9,2.5,'移动端响应式适配',['≤768px全局CSS媒体查询','表格12px/按钮缩小padding','弹窗92%宽/抽屉220px','分页居中/表单项间距缩小'],E)
code(s,M,4,6.5,1.8,'评价回复处理',['@PutMapping("/{id}/reply")','reviewService.reply(id, body.get("reply"))','// ReviewServiceImpl:','review.setReply(reply)','review.setReplyTime(LocalDateTime.now())','reviewMapper.updateById(review)'])
code(s,7,4,5.8,1.8,'移动端CSS响应式(全局)',['@media (max-width: 768px) {','  .el-table { font-size: 12px }','  .el-button--small { padding: 4px 8px }','  .el-pagination { justify-content: center }','  .el-dialog { width: 92%!important }','  .el-drawer { width: 220px!important }','}'])
img_placeholder(s,4.7,4.3,3.9,2,'[此处插入: Banner管理页面截图]')

# ====== P20 AI客服创新(1) ======
s=sl()
hdr(s,'07 AI 智能客服 (核心亮点)','技术方案与工作流程','DeepSeek大语言模型 + RAG检索增强生成 + 关联词映射引擎')
card(s,M,1.2,5.8,2.3,'技术方案',['模型: deepseek-chat (RestTemplate调用)','System Prompt: 动态构建(含商品列表+规则)','  禁止Markdown符号, 口语化回复','  角色: QJ商城AI助手小Q, 工号26793','关联词引擎: HashMap<String,String[]>','  11个品类 60+关键词映射','推荐协议: [RECOMMEND]id1,id2,id3[/RECOMMEND]'],P)
# 工作流程(竖向)
bx(s,7,1.2,5.8,2.8,W)
t=tx(s,7.15,1.3,5.5,0.4);S(t,'6步工作流程',fs=16,b=True,c=D)
steps=['1. 用户发送消息(支持泛指"我想要饮料")','  ↓','2. 提取关键词+关联词扩展(11品类→DB搜索词)','  ↓','3. LambdaQueryWrapper搜索DB(上架+有库存+销量Top10)','  ↓','4. 构建System Prompt(商品列表+历史10条)+调用DeepSeek','  ↓','5. API返回→解析[RECOMMEND]标记→上限3个商品卡片','  ↓','6. API失败→catch→getFallbackReply()降级预设回复(<1s)']
for ln in steps:
    if ln.startswith('  ↓'):A(t,ln,fs=10,c=G,sb=2)
    else:A(t,ln,fs=10,c=D,sb=5)
code(s,M,3.8,12.2,2.8,'关联词映射引擎 (11品类60+关键词)',['static final Map<String,String[]> RELATED_WORDS = new LinkedHashMap<>()','{{"饮料",{"气泡水","可乐","橙汁","果汁","水","茶","咖啡"}}};','{{"手机",{"iPhone","华为","小米","OPPO","vivo","三星"}}};','{{"电脑",{"MacBook","ThinkPad","MateBook","笔记本"}}};','{{"衣服",{"羽绒","卫衣","大衣","CK","优衣库","Nike","ZARA"}}};','{{"护肤",{"SK-II","雅诗兰黛","海蓝之谜","神仙水","精华"}}};','{{"家电",{"冰箱","洗衣机","空调","电视","戴森","吸尘器"}}};','// 共11品类: 饮料/手机/电脑/耳机/零食/衣服/鞋/护肤/家电/酒/手表'])

# ====== P21 AI客服创新(2) ======
s=sl()
hdr(s,'07 AI 智能客服 (核心亮点)','前端展示与效果示例','聊天气泡 + 商品推荐卡片 + 降级策略')
card(s,M,1.2,5.8,2.5,'前端聊天UI (CustomerService.vue)',['欢迎区: 机器人emoji+工号+快捷问题(5个标签)','消息气泡: 用户蓝色右对齐 / AI 白色左对齐','AI思考: 三点闪烁动画(blink keyframes)','推荐卡片: 56x44图片+名称+价格+跳转箭头','  @click="goProduct(p.id)"→/client/product/:id','Markdown过滤: /[*#_`~>-]/g正则清理'],P)
card(s,7,1.2,5.8,2.5,'效果演示',['用户: "我想要饮料"','AI: "为您推荐以下饮料哦~"','【元气森林气泡水】无糖低卡 ¥59.90','【智利进口车厘子】JJ级 ¥199.00','','用户: "推荐一款手机"','AI: 推荐iPhone15ProMax/小米14Pro/三星S24','各附价格+购买链接(点击卡片跳转)'],O)
code(s,M,4,6.5,1.6,'AI推荐解析(后端)',['// 解析[RECOMMEND]id1,id2,id3[/RECOMMEND]','if (content.contains("[RECOMMEND]")) {','  String ids = content.substring(start+11,end)','  for (String idStr : ids.split(",")) {','    Long id = Long.parseLong(idStr.trim())','    // 匹配推荐商品→组装卡片数据','    // {id, name, price, image, link}','  }','}'])
code(s,7,4,5.8,1.6,'降级策略',['} catch (Exception e) {','  return fallbackReply(userMessage)','  // 预设关键词回复:','  // "订单"→发货时间/查看指引','  // "退货"→退款流程说明','  // "运费"→包邮政策','  // 默认→引导性回复','}'])
img_placeholder(s,M,5.9,5.8,1,'[此处插入: AI客服对话界面截图 — 展示推荐卡片效果]')
img_placeholder(s,7,5.9,5.8,1,'[此处插入: 管理端客服管理界面截图]')

# ====== P22 作品测试 ======
s=sl()
hdr(s,'08 作品测试','测试概览','37个测试用例 | 100%通过率 | 双端全覆盖')
# 顶部统计
for i,(v,lb,c) in enumerate([('37','测试用例总数',P),('100%','测试通过率',E),('24','用户端用例',D),('13','管理端用例',O),('Chrome/Edge\nSafari/微信','兼容性验证',P)]):
    stat_card(s,M+i*2.35,1.2,2.2,1.5,v,lb,c)
# 中部: 测试环境
bx(s,M,3,5.8,3.2,W)
t=tx(s,M+0.3,3.1,5.2,2.8)
S(t,'测试环境',fs=18,b=True,c=D)
for ln in ['服务器: 阿里云ECS 2核4G CentOS7.9','数据库: MySQL 8.0 (数据库qj_shop)','浏览器: Chrome120+/Edge120+/Safari17+/微信','设备: iPhone14Pro / AndroidChrome / 鸿蒙模拟器','PC分辨率: 1920x1080 / 1366x768','数据: 40商品/12订单/24评价/8用户']:
    A(t,ln,fs=13,c=G,sb=8)
# 右: 测试用例
bx(s,7,3,5.8,3.2,W)
t=tx(s,7.3,3.1,5.2,2.8)
S(t,'用户端测试用例 (24个) 全部通过',fs=16,b=True,c=E)
for ln in ['注册(合法/重复) | 登录(正确/错误)','首页加载 | 搜索 | 分类浏览','商品详情(移动/PC双端) | 规格选择','购物车(加入/修改/删除/全选/结算)','下单(地址→确认→支付) | 订单(Tab/取消)','AI客服(泛指推荐/口语回复/降级)','个人中心(统计/头像拍照) | 地址(GPS)']:
    A(t,ln,fs=11,c=G,sb=6)
# 底部
bx(s,M,6.5,11.7,0.5,L)
t=tx(s,M+0.3,6.52,11.1,0.45)
S(t,'管理端13个用例全部通过: 仪表盘 | 用户CRUD | 分类CRUD | 商品三级 | 订单发货 | Banner上传 | 评价回复 | 移动适配 | Excel导出',fs=12,b=True,c=P,a=PP_ALIGN.CENTER)

# ====== P23 兼容性&性能 ======
s=sl()
hdr(s,'08 作品测试','兼容性与性能验证','多浏览器/多设备/多分辨率全覆盖')
card(s,M,1.2,5.8,2.8,'浏览器与设备兼容性 (全部通过)',['Chrome 120+ ✓  Edge 120+ ✓  Safari 17+ ✓','微信内置浏览器(Android) ✓  (iOS) ✓','鸿蒙DevEco模拟器: WebView正常加载 ✓','响应式: 320px~1920px全覆盖 ✓','  小屏(<500px)商品3列 | 中屏4列 | 大屏5列','高分辨率手机: UA检测→正确显示移动布局 ✓'],P)
card(s,7,1.2,5.8,2.8,'性能指标',['页面首次加载: <3秒 ✓','商品列表查询: <1秒 ✓','订单提交: <2秒 ✓','AI客服(正常): 3-5秒 ✓','AI客服(降级): <1秒 ✓','并发50用户浏览: 无明显延迟 ✓','JAR启动: ~7秒 | JAR大小: 64MB','前端dist: 38文件 | 后端class: 80+文件'],D)
code(s,M,4.3,5.8,1.8,'Druid连接池监控',['// /druid/index.html 实时监控','// 活跃连接数 | SQL执行时间','// 慢SQL记录(>2s) | 连接池状态'])
code(s,7,4.3,5.8,1.8,'性能优化建议',['// 可引入Redis缓存热门数据','// 可引入RabbitMQ异步处理订单','// 图片可迁移至OSS对象存储','// 可配置CDN加速静态资源'])
img_placeholder(s,M,6.4,5.8,0.6,'[此处插入: 浏览器兼容性测试截图 | 移动端响应式展示]')
img_placeholder(s,7,6.4,5.8,0.6,'[此处插入: Druid监控面板截图 | 性能测试数据]')

# ====== P24 技术问题(1) ======
s=sl()
hdr(s,'09 技术问题与解决','8个技术问题全部攻克 (Page 1/2)','开发过程中遇到的关键技术难点及解决方案')
p1=[('Vant4组件名变更','GoodsAction→ActionBar','查阅CHANGELOG迁移'),
    ('高分辨率手机误判PC','1080px手机显示PC布局','UA+触摸+宽度三维检测'),
    ('Vite代理rewrite','/api/路径不匹配Controller','去rewrite+统一路径'),
    ('Banner时间过滤NULL','startTime NULL→不命中','SQL增加IS NULL OR判断')]
for i,(name,prob,sol) in enumerate(p1):
    row=i//2;col=i%2
    lx=M+col*6.15;ly=1.2+row*2.7
    bx(s,lx,ly,5.85,2.4,W)
    bx(s,lx,ly,5.85,0.05,O if i<2 else O)
    t=tx(s,lx+0.2,ly+0.12,5.45,2.1)
    S(t,f'问题{i+1}: {name}',fs=15,b=True,c=D)
    A(t,f'现象: {prob}',fs=12,c=G,sb=6)
    A(t,f'解决: {sol}',fs=12,c=E,sb=6)

# ====== P25 技术问题(2) ======
s=sl()
hdr(s,'09 技术问题与解决','8个技术问题全部攻克 (Page 2/2)','项目经验总结')
p2=[('Linux MySQL Socket','localhost→socket连接失败','JDBC URL改为127.0.0.1'),
    ('SPA路由刷新500','/admin直接访问报500','SpaFilter最高优先级转发index.html'),
    ('文件上传boundary','手动设Content-Type缺boundary','让axios自动生成请求头'),
    ('注册密码双重BCrypt','Controller+Service各自加密','去掉Controller重复编码')]
for i,(name,prob,sol) in enumerate(p2):
    row=i//2;col=i%2
    lx=M+col*6.15;ly=1.2+row*2.7
    bx(s,lx,ly,5.85,2.4,W)
    bx(s,lx,ly,5.85,0.05,O)
    t=tx(s,lx+0.2,ly+0.12,5.45,2.1)
    S(t,f'问题{i+5}: {name}',fs=15,b=True,c=D)
    A(t,f'现象: {prob}',fs=12,c=G,sb=6)
    A(t,f'解决: {sol}',fs=12,c=E,sb=6)
# 经验总结
bx(s,M,5.2,11.7,1.5,L)
t=tx(s,M+0.3,5.3,11.1,1.2)
S(t,'项目经验总结',fs=16,b=True,c=D)
for ln in ['文档先行: 升级依赖前查阅CHANGELOG和Migration Guide','编码规范: 密码加密统一在Service层, Controller只做参数接收','防御编程: SQL查询考虑NULL边界条件 | JDBC优先TCP连接避免Socket路径问题','代理配置: 理解前端代理路径映射关系 | SPA需配置过滤器回退index.html']:
    A(t,ln,fs=12,c=G,sb=8)

# ====== P26 项目成果 ======
s=sl()
hdr(s,'10 项目总结','成果汇总','80+源文件 | 19 Controller | 27 Vue页面 | 3终端 | 12张表 | 37测试 | 25+提交')
for i,(v,lb,desc,c) in enumerate([
    ('80+','源文件','Java+Vue+ArkTS+SQL',P),('19','Controller','5公共+5用户+9管理',D),('11','Service','业务逻辑完整实现',E),
    ('12','数据表','全部逻辑删除+快照',O),('27','Vue页面','15用户+12管理+双端',P),('3','终端','Web+PC+鸿蒙',D),
    ('37','测试用例','100%通过率',E),('25+','Git提交','完整开发日志',O),('64MB','FatJAR','前后端一体化部署',P)]):
    row=i//3;col=i%3
    stat_card(s,M+col*4.1,1.2+row*2.1,3.9,1.8,v,lb,c)
    t=tx(s,M+col*4.1+0.1,1.2+row*2.1+1.5,3.7,0.3)
    S(t,desc,fs=10,c=G,a=PP_ALIGN.CENTER)

# ====== P27 创新点 ======
s=sl()
hdr(s,'10 项目总结','核心创新点','五大创新 + 技术难度/实用价值星级评估')
for i,(title,desc,c) in enumerate([
    ('AI客服 + RAG简化实现','DeepSeek大模型+关联词引擎(11品类60+词)\n自然语言→商品推荐→聊天卡片→购买',P),
    ('三维设备检测策略','UA+触摸屏+宽度三维判定\n解决高分辨率手机误判PC端问题',D),
    ('GPS/IP双重定位','W3C Geolocation→PermissionDenied\n→Nominatim→ipapi→ipsb回退',E),
    ('前后端一体化部署','Vue dist→SpringBoot static→Maven Package\nFatJAR 64MB, java -jar一键启动',O),
    ('12张表完整数据模型','订单快照机制+逻辑删除+雪花ID\n覆盖商品/订单/用户/内容全链路',P)]):
    row=i//2;col=i%2
    lx=M+col*6.15;ly=1.2+row*1.8
    bx(s,lx,ly,5.85,1.55,W)
    bx(s,lx,ly,5.85,0.05,c)
    t=tx(s,lx+0.15,ly+0.12,5.55,1.3)
    S(t,title,fs=15,b=True,c=c)
    for ln in desc.split('\n'):A(t,ln,fs=12,c=G,sb=3)
    A(t,'技术难度: ★★★★ | 实用价值: ★★★★★',fs=10,c=O,sb=4)

# ====== P28 不足与改进 ======
s=sl()
hdr(s,'10 项目总结','不足与改进方向','六维度持续优化路径')
for i,(title,lines,c) in enumerate([
    ('支付真实化',['当前: 模拟支付→状态流转','改进: 接入支付宝/微信支付SDK','→ 真实回调+签名验证'],P),
    ('AI推荐精准化',['当前: 关联词映射→DB搜索','改进: 向量数据库语义搜索','→ 协同过滤+深度学习推荐'],D),
    ('图片存储优化',['当前: base64存DB longtext','改进: 迁移OSS对象存储','→ URL引用, 按需加载'],E),
    ('性能优化',['当前: 单体应用 单机部署','改进: Redis缓存+MQ异步','→ SpringCloud微服务化'],O),
    ('功能扩展',['当前: 覆盖电商核心流程','改进: 优惠券/秒杀/物流','→ ECharts数据可视化'],P),
    ('安全增强',['当前: BCrypt+JWT+Auth','改进: HTTPS+验证码防刷','→ 接口限流+审计日志'],D)]):
    row=i//3;col=i%3
    lx=M+col*4.1;ly=1.2+row*3
    bx(s,lx,ly,3.9,2.7,W)
    bx(s,lx,ly,3.9,0.05,c)
    t=tx(s,lx+0.15,ly+0.12,3.6,2.3)
    S(t,title,fs=15,b=True,c=c)
    for ln in lines:A(t,ln,fs=12,c=G,sb=6)

# ====== P29 结束页 ======
s=sl()
bx(s,0,0,13.333,7.5,P)
t=tx(s,2,2.2,9.3,2)
S(t,'感谢聆听',fs=52,b=True,c=W,a=PP_ALIGN.CENTER)
A(t,'恳请各位老师批评指正',fs=24,c=W,a=PP_ALIGN.CENTER,sb=20)
bx(s,3.5,4.5,6.3,0.01,W)
t=tx(s,3,5,7.3,1.5)
S(t,'"柒玖商店" — 基于SpringBoot多前端多平台的商城系统',fs=14,c=W,a=PP_ALIGN.CENTER)
A(t,'李 鑫 | 202330302152 | 软件技术 2307 班 | 彭德宇 老师 | 2026年5月',fs=12,c=W,a=PP_ALIGN.CENTER,sb=8)

# ====== 保存 ======
path=r"D:\QJShop\柒玖商店_毕业设计答辩_优化版.pptx"
prs.save(path)
print(f"SAVED: {path}")
print(f"SLIDES: {len(prs.slides)}")
